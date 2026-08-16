"""Chart handling module for EasyPPTX."""

from __future__ import annotations

from typing import TYPE_CHECKING, Any, ClassVar

from pptx.chart.chart import Chart as PPTXChart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

from easypptx.common import is_dataframe
from easypptx.positioning import PositionType

if TYPE_CHECKING:
    import pandas as pd

    from easypptx.slide import Slide


def extract_categories_values(
    data: Any,
    category_column: str | int | None = None,
    value_columns: str | list[str] | int | list[int] | None = None,
) -> tuple[list, list]:
    """Extract category and value lists from tabular chart data.

    Args:
        data: A pandas DataFrame, or a list of lists whose first row is a header
        category_column: Name or index of the category column (default: first column)
        value_columns: Name(s) or index(es) of the value column(s); only the
            first is used (default: second column)

    Returns:
        A (categories, values) tuple of lists

    Raises:
        ValueError: If a named column is missing or the data has too few columns
    """
    first_value = value_columns[0] if isinstance(value_columns, list) else value_columns

    if is_dataframe(data):
        if category_column is None:
            category_column = data.columns[0]
        if first_value is None:
            if len(data.columns) < 2:
                raise ValueError("DataFrame must have at least two columns for automatic value extraction")
            first_value = data.columns[1]

        def df_column(column: Any, kind: str) -> Any:
            # Prefer label lookup; fall back to positional access for ints
            if column in data.columns:
                return data[column]
            if isinstance(column, int) and 0 <= column < len(data.columns):
                return data.iloc[:, column]
            raise ValueError(f"{kind} column '{column}' not found in DataFrame")

        return df_column(category_column, "Category").tolist(), df_column(first_value, "Value").tolist()

    if not data or len(data) < 2:
        return [], []

    header = data[0]

    def column_index(column: str | int | None, default: int, kind: str) -> int:
        if column is None:
            return default
        if isinstance(column, int):
            return column
        try:
            return header.index(column)
        except ValueError:
            raise ValueError(f"{kind} column '{column}' not found in header") from None

    cat_idx = column_index(category_column, 0, "Category")
    if first_value is None and len(header) < 2:
        raise ValueError("Data must have at least two columns for automatic value extraction")
    val_idx = column_index(first_value, 1, "Value")

    categories = [row[cat_idx] for row in data[1:]]
    values = [row[val_idx] for row in data[1:]]
    return categories, values


def extract_chart_series(
    data: Any,
    category_column: str | int | None = None,
    value_columns: str | list[str] | int | list[int] | None = None,
) -> tuple[list, dict[str, list]]:
    """Extract categories and one or more named value series from tabular data.

    Args:
        data: A pandas DataFrame, or a list of lists whose first row is a header
        category_column: Name or index of the category column (default: first column)
        value_columns: Name(s) or index(es) of the value column(s)
            (default: the second column)

    Returns:
        (categories, series) where series maps series name -> list of values

    Raises:
        ValueError: If a named column is missing or the data has too few columns
    """
    columns = value_columns if isinstance(value_columns, list) else [value_columns]

    if is_dataframe(data):
        if category_column is None:
            category_column = data.columns[0]

        def df_column(column: Any, kind: str) -> Any:
            if column in data.columns:
                return data[column]
            if isinstance(column, int) and 0 <= column < len(data.columns):
                return data.iloc[:, column]
            raise ValueError(f"{kind} column '{column}' not found in DataFrame")

        categories = df_column(category_column, "Category").tolist()
        series: dict[str, list] = {}
        for column in columns:
            if column is None:
                if len(data.columns) < 2:
                    raise ValueError("DataFrame must have at least two columns for automatic value extraction")
                column = data.columns[1]
            values = df_column(column, "Value")
            series[str(getattr(values, "name", column))] = values.tolist()
        return categories, series

    if not data or len(data) < 2:
        return [], {}

    header = data[0]

    def column_index(column: str | int | None, default: int, kind: str) -> int:
        if column is None:
            return default
        if isinstance(column, int):
            return column
        try:
            return header.index(column)
        except ValueError:
            raise ValueError(f"{kind} column '{column}' not found in header") from None

    cat_idx = column_index(category_column, 0, "Category")
    categories = [row[cat_idx] for row in data[1:]]
    series = {}
    for column in columns:
        if column is None and len(header) < 2:
            raise ValueError("Data must have at least two columns for automatic value extraction")
        val_idx = column_index(column, 1, "Value")
        series[str(header[val_idx])] = [row[val_idx] for row in data[1:]]
    return categories, series


class Chart:
    """Class for handling chart operations in PowerPoint slides.

    This class provides methods for creating and manipulating charts on slides.

    Examples:
        ```python
        # Create a chart object
        chart = Chart(slide)

        # Add a bar chart
        chart.add_bar(
            categories=["A", "B", "C"],
            values=[1, 2, 3],
            title="Sample Bar Chart"
        )

        # Add a pie chart from DataFrame
        import pandas as pd
        df = pd.DataFrame({"Category": ["A", "B", "C"], "Value": [10, 20, 30]})
        chart.from_dataframe(df, chart_type="pie", x=2, y=2)
        ```
    """

    CHART_TYPES: ClassVar = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }

    LEGEND_POSITIONS: ClassVar = {
        "right": XL_LEGEND_POSITION.RIGHT,
        "left": XL_LEGEND_POSITION.LEFT,
        "top": XL_LEGEND_POSITION.TOP,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "corner": XL_LEGEND_POSITION.CORNER,
    }

    def __init__(self, slide_obj: Slide) -> None:
        """Initialize a Chart object.

        Args:
            slide_obj: The Slide object to add charts to
        """
        self.slide = slide_obj

    def add(
        self,
        chart_type: str,
        categories: list,
        values: list | None = None,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        title: str | None = None,
        has_legend: bool = True,
        series: dict[str, list] | None = None,
        show_values: bool = False,
        number_format: str | None = None,
        x_title: str | None = None,
        y_title: str | None = None,
        y_min: float | None = None,
        y_max: float | None = None,
        palette: list | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a chart to the slide.

        Args:
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', 'area', 'scatter')
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            has_legend: Whether to show legend (default: True)
            series: Mapping of series name -> values for multi-series charts (default: None)
            show_values: Draw data labels on the series (default: False)
            number_format: Excel-style number format for data labels,
                e.g. "#,##0" or "0.0%" (default: None)
            x_title: Category-axis title (default: None)
            y_title: Value-axis title (default: None)
            y_min: Lower value-axis limit (default: None)
            y_max: Upper value-axis limit (default: None)
            palette: Series colors as color names or RGB tuples (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object

        Raises:
            ValueError: If chart_type is not supported or data is invalid
        """
        if chart_type not in self.CHART_TYPES:
            raise ValueError(
                f"Unsupported chart type: {chart_type}. Supported types: {', '.join(self.CHART_TYPES.keys())}"
            )

        if series is None:
            if values is None:
                raise ValueError("Provide either 'values' or 'series'")
            series = {"Series 1": values}

        for name, series_values in series.items():
            if len(categories) != len(series_values):
                raise ValueError(f"Series '{name}' must have the same length as categories")

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, series_values in series.items():
            chart_data.add_series(name, series_values)

        # Get slide dimensions for percentage conversion
        slide_width = self.slide._get_slide_width()
        slide_height = self.slide._get_slide_height()

        # Convert position values to inches
        x_inches = self.slide._convert_position(x, slide_width)
        y_inches = self.slide._convert_position(y, slide_height)
        width_inches = self.slide._convert_position(width, slide_width)
        height_inches = self.slide._convert_position(height, slide_height)

        chart_shape = self.slide.pptx_slide.shapes.add_chart(
            self.CHART_TYPES[chart_type],
            Inches(x_inches),
            Inches(y_inches),
            Inches(width_inches),
            Inches(height_inches),
            chart_data,  # ty: ignore[invalid-argument-type]
        )

        chart = chart_shape.chart  # ty: ignore[unresolved-attribute]

        # Set chart title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False

        # Set legend visibility
        chart.has_legend = has_legend
        if has_legend:
            legend_position = kwargs.get("legend_position", "right")
            if isinstance(legend_position, str) and legend_position in self.LEGEND_POSITIONS:
                chart.legend.position = self.LEGEND_POSITIONS[legend_position]

        # Data labels
        if show_values or number_format:
            plot = chart.plots[0]
            plot.has_data_labels = True
            if number_format:
                plot.data_labels.number_format = number_format
                plot.data_labels.number_format_is_linked = False

        # Axis titles and value-axis limits (not every chart type has axes)
        chart_any: Any = chart
        if chart_type != "pie" and (x_title or y_title or y_min is not None or y_max is not None):
            import warnings

            try:
                if x_title:
                    axis = chart_any.category_axis
                    axis.has_title = True
                    axis.axis_title.text_frame.text = x_title
                if y_title:
                    axis = chart_any.value_axis
                    axis.has_title = True
                    axis.axis_title.text_frame.text = y_title
                if y_min is not None:
                    chart_any.value_axis.minimum_scale = y_min
                if y_max is not None:
                    chart_any.value_axis.maximum_scale = y_max
            except ValueError as err:
                warnings.warn(f"Axis options not supported for chart type {chart_type!r}: {err}", stacklevel=2)

        # Series colors from an explicit palette (e.g. the deck theme)
        if palette:
            from easypptx.common import resolve_color

            for i, chart_series in enumerate(chart.plots[0].series):
                rgb = resolve_color(palette[i % len(palette)])
                if rgb is not None:
                    chart_series.format.fill.solid()
                    chart_series.format.fill.fore_color.rgb = rgb

        return chart

    def add_bar(
        self,
        categories: list,
        values: list,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a bar chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="bar",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def add_column(
        self,
        categories: list,
        values: list,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a column chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="column",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def add_pie(
        self,
        categories: list,
        values: list,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a pie chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="pie",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def from_dataframe(
        self,
        df: pd.DataFrame,
        chart_type: str,
        category_column: str,
        value_column: str,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        title: str | None = None,
        has_legend: bool = True,
        **kwargs: Any,
    ) -> PPTXChart:
        """Create a chart from a pandas DataFrame.

        Args:
            df: Pandas DataFrame
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', 'area', 'scatter')
            category_column: Column name to use for categories
            value_column: Column name to use for values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            has_legend: Whether to show legend (default: True)
            series: Mapping of series name -> values for multi-series charts (default: None)
            show_values: Draw data labels on the series (default: False)
            number_format: Excel-style number format for data labels,
                e.g. "#,##0" or "0.0%" (default: None)
            x_title: Category-axis title (default: None)
            y_title: Value-axis title (default: None)
            y_min: Lower value-axis limit (default: None)
            y_max: Upper value-axis limit (default: None)
            palette: Series colors as color names or RGB tuples (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object

        Raises:
            ValueError: If columns don't exist in DataFrame
        """
        if category_column not in df.columns:
            raise ValueError(f"Category column '{category_column}' not found in DataFrame")
        if value_column not in df.columns:
            raise ValueError(f"Value column '{value_column}' not found in DataFrame")

        categories = df[category_column].tolist()
        values = df[value_column].tolist()

        return self.add(
            chart_type=chart_type,
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            has_legend=has_legend,
            **kwargs,
        )
