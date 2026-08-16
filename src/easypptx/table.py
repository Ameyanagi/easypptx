"""Table handling module for EasyPPTX."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.table import Table as PPTXTable
from pptx.util import Inches, Pt

from easypptx.positioning import PositionType

# Built-in PowerPoint table style GUIDs, addressable by small integer ids.
# Sources: the standard Office table style gallery.
TABLE_STYLE_GUIDS = {
    0: "{2D5ABB26-0587-4C30-8999-92F81FD0307C}",  # No Style, No Grid
    1: "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",  # Medium Style 2 - Accent 1
    2: "{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}",  # Medium Style 2 - Accent 2
    3: "{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}",  # Medium Style 2 - Accent 3
    4: "{00A15C55-8517-42AA-B614-E9B94910E393}",  # Medium Style 2 - Accent 4
    5: "{7DF18680-E054-41AD-8BC1-D1AEF772440D}",  # Medium Style 2 - Accent 5
    6: "{93296810-A885-4BE3-A3E7-6D5BEEA58F35}",  # Medium Style 2 - Accent 6
    7: "{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}",  # Light Style 1
    8: "{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}",  # Light Style 1 - Accent 1
}


def _apply_table_style(table: PPTXTable, style: int | str) -> None:
    """Set the table's built-in style by small integer id or literal GUID."""
    from pptx.oxml.ns import qn

    guid = TABLE_STYLE_GUIDS.get(style) if isinstance(style, int) else style
    if guid is None:
        raise ValueError(f"Unknown table style id: {style!r}. Known ids: {sorted(TABLE_STYLE_GUIDS)} or a GUID string")

    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = tbl.makeelement(qn("a:tblPr"), {})
        tbl.insert(0, tblPr)
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = tblPr.makeelement(qn("a:tableStyleId"), {})
        tblPr.append(style_id)
    style_id.text = guid


if TYPE_CHECKING:
    import pandas as pd

    from easypptx.slide import Slide


def _is_finite_number(value: object) -> bool:
    """True for real, finite, non-bool numbers (incl. numpy scalars)."""
    import math
    import numbers

    if isinstance(value, bool) or not isinstance(value, numbers.Real):
        return False
    try:
        return math.isfinite(float(value))
    except (TypeError, ValueError, OverflowError):
        return False


def apply_number_format(rows: list[list], number_format: str | dict, has_header: bool = True) -> list[list]:
    """Format numeric body cells with Python format specs.

    Args:
        rows: Table rows
        number_format: A format string like "{:,.1f}" applied to every
            numeric cell, or a dict mapping column name/index to a format
            (column names win over positional indexes)
        has_header: Whether the first row is a header and stays unformatted
            (default: True)

    Returns:
        A new rows list with formatted string values
    """
    if not rows:
        return rows
    header = rows[0] if has_header else []
    body_start = 1 if has_header else 0

    def format_for(col: int) -> str | None:
        if isinstance(number_format, str):
            return number_format
        # Column names take precedence; a key matching a header name is
        # never reused as a positional index
        name = header[col] if col < len(header) else None
        if name is not None and name in number_format:
            return number_format[name]
        if col in number_format and col not in header:
            return number_format[col]
        return None

    formatted = [list(row) for row in rows[:body_start]]
    for row in rows[body_start:]:
        new_row = []
        for col, value in enumerate(row):
            fmt = format_for(col)
            if fmt is not None and _is_finite_number(value):
                new_row.append(fmt.format(value))
            else:
                new_row.append(value)
        formatted.append(new_row)
    return formatted


def shade_cells_by_value(
    table: PPTXTable,
    rows: list[list],
    shade_columns: list,
    shade_color: str | tuple[int, int, int] = "blue",
    has_header: bool = True,
) -> None:
    """Tint body-cell backgrounds by value: min stays white, max gets shade_color.

    NaN/infinite values and non-numbers are left unshaded.

    Args:
        table: The rendered python-pptx table
        rows: The table data
        shade_columns: Column names or positional indexes to shade
            (names win over indexes)
        shade_color: The full-intensity tint color
        has_header: Whether the first row is a header and stays unshaded
            (default: True)
    """
    from pptx.dml.color import RGBColor

    from easypptx.common import resolve_color

    rgb = resolve_color(shade_color)
    body_start = 1 if has_header else 0
    if rgb is None or len(rows) <= body_start:
        return
    header = rows[0] if has_header else []
    n_cols = max(len(row) for row in rows)

    def column_index(column: object) -> int | None:
        # Column names take precedence over positional indexes
        if column in header:
            return header.index(column)
        if isinstance(column, int) and 0 <= column < n_cols:
            return column
        return None

    for column in shade_columns:
        col = column_index(column)
        if col is None:
            raise ValueError(f"shade column '{column}' not found in header")
        values = [float(row[col]) for row in rows[body_start:] if _is_finite_number(row[col])]
        if not values:
            continue
        low, high = min(values), max(values)
        span = (high - low) or 1.0
        for i, row in enumerate(rows[body_start:], start=body_start):
            value = row[col]
            if not _is_finite_number(value):
                continue
            t = (float(value) - low) / span
            blended = RGBColor(
                round(255 + (rgb[0] - 255) * t),
                round(255 + (rgb[1] - 255) * t),
                round(255 + (rgb[2] - 255) * t),
            )
            cell = table.cell(i, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = blended


def _looks_numeric(text: str) -> bool:
    """True for cell text that reads as a number (with separators/percent)."""
    cleaned = text.strip().replace(",", "").replace("%", "").replace("\u00a5", "").replace("$", "")
    if not cleaned:
        return False
    try:
        float(cleaned)
    except ValueError:
        return False
    return True


def apply_table_theme(table: PPTXTable, spec: dict, has_header: bool = True) -> None:
    """Apply a theme's table styling: header fill, banding, alignment, fonts.

    Args:
        table: The rendered python-pptx table
        spec: Theme table spec with optional keys header_fill, header_color,
            band_fills (list of two fills), body_color, font_name,
            font_size, header_font_size
        has_header: Whether row 0 is a header row (default: True)
    """
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    from easypptx.common import resolve_color

    header_fill = resolve_color(spec.get("header_fill"))
    header_color = resolve_color(spec.get("header_color"))
    band_fills = [resolve_color(c) for c in spec.get("band_fills", [])]
    body_color = resolve_color(spec.get("body_color"))
    font_name = spec.get("font_name")
    font_size = spec.get("font_size")
    header_font_size = spec.get("header_font_size", font_size)

    n_rows = len(table.rows)
    n_cols = len(table.columns)
    body_start = 1 if has_header else 0

    for r in range(n_rows):
        is_header = has_header and r == 0
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

            # Fills: header color, then subtle banding
            if is_header and header_fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif not is_header and band_fills and all(band_fills):
                cell.fill.solid()
                cell.fill.fore_color.rgb = band_fills[(r - body_start) % len(band_fills)]

            for paragraph in cell.text_frame.paragraphs:
                if font_name:
                    paragraph.font.name = font_name
                size = header_font_size if is_header else font_size
                if size:
                    paragraph.font.size = Pt(size)
                color = header_color if is_header else body_color
                if color is not None:
                    paragraph.font.color.rgb = color
                # Numbers right-align for easy scanning; header follows body
                if (not is_header and _looks_numeric(cell.text)) or (
                    is_header and n_rows > body_start and _looks_numeric(table.cell(body_start, c).text)
                ):
                    paragraph.alignment = PP_ALIGN.RIGHT


class Table:
    """Class for handling table operations in PowerPoint slides.

    This class provides methods for creating and manipulating tables on slides.

    Examples:
        ```python
        # Create a table object
        table = Table(slide)

        # Add a simple table
        table.add([["Header 1", "Header 2"], ["Value 1", "Value 2"]])

        # Add a table from pandas DataFrame
        import pandas as pd
        df = pd.DataFrame({'A': [1, 2], 'B': [3, 4]})
        table.from_dataframe(df, x=2, y=2)
        ```
    """

    def __init__(self, slide_obj: Slide) -> None:
        """Initialize a Table object.

        Args:
            slide_obj: The Slide object to add tables to
        """
        self.slide = slide_obj

    def add(
        self,
        data: list,
        x: PositionType = 5,
        y: PositionType = 20,
        width: PositionType | None = None,
        height: PositionType | None = None,
        first_row_header: bool = True,
        style: int | str | None = None,
        has_header: bool | None = None,
    ) -> PPTXTable:
        """Add a table to the slide.

        Args:
            data: 2D list of table data
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Total width in inches or percentage (default: None, auto-sized)
            height: Total height in inches or percentage (default: None, auto-sized)
            first_row_header: Whether to format the first row as a header (default: True)
            style: Table style ID (default: None)
            has_header: Alias for first_row_header; takes precedence when given (default: None)

        Returns:
            The created table object
        """
        if has_header is not None:
            first_row_header = has_header
        if not data:
            raise ValueError("Table data cannot be empty")

        rows = len(data)
        cols = len(data[0])

        # Ensure all rows have the same number of columns
        for row in data:
            if len(row) != cols:
                raise ValueError("All rows must have the same number of columns")

        # Get slide dimensions for percentage conversion
        slide_width = self.slide._get_slide_width()
        slide_height = self.slide._get_slide_height()

        # Convert position values to inches
        x_inches = self.slide._convert_position(x, slide_width)
        y_inches = self.slide._convert_position(y, slide_height)

        # Create table shape
        # Default width based on columns if None, otherwise convert from position
        width_inches = cols * 2.0 if width is None else self.slide._convert_position(width, slide_width)

        # Default height based on rows if None, otherwise convert from position
        height_inches = rows * 0.5 if height is None else self.slide._convert_position(height, slide_height)

        table_shape = self.slide.pptx_slide.shapes.add_table(
            rows, cols, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )
        table = table_shape.table

        # Fill table data
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)

                # Format header row
                if first_row_header and i == 0:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(14)

        # Apply table style if specified
        if style is not None:
            _apply_table_style(table, style)

        return table

    def from_dataframe(
        self,
        df: pd.DataFrame,
        x: PositionType = 5,
        y: PositionType = 20,
        width: PositionType | None = None,
        height: PositionType | None = None,
        include_index: bool = False,
        first_row_header: bool = True,
        style: int | None = None,
    ) -> PPTXTable:
        """Add a table from a pandas DataFrame.

        Args:
            df: Pandas DataFrame
            x: X position in inches (default: 1.0)
            y: Y position in inches (default: 1.0)
            width: Total width in inches (default: None, auto-sized)
            height: Total height in inches (default: None, auto-sized)
            include_index: Whether to include DataFrame index (default: False)
            first_row_header: Whether to format column names as headers (default: True)
            style: Table style ID (default: None)

        Returns:
            The created table object
        """
        # Convert DataFrame to list format
        if include_index:
            # The header row needs a leading cell for the index column
            data = [[df.index.name or "", *list(df.columns)]]
            for idx, row in df.iterrows():
                data.append([str(idx), *list(row)])
        else:
            data = [list(df.columns), *df.values.tolist()]

        return self.add(
            data=data,
            x=x,
            y=y,
            width=width,
            height=height,
            first_row_header=first_row_header,
            style=style,
        )
