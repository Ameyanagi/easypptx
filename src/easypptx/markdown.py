"""Markdown-to-presentation conversion for EasyPPTX.

Converts a markdown document into a deck:

- optional frontmatter (``--- ... ---``) with ``theme``, ``template``,
  ``aspect_ratio``, and ``subtitle`` settings
- ``# Heading`` becomes the title slide, ``## Heading`` starts a new slide
- ``---`` on its own line forces a slide break
- ``-`` / ``*`` / ``1.`` lists become bullets (2-space indent nests)
- ``![alt](path)`` images, GFM ``|`` tables, and fenced code blocks
- ``<!-- notes: ... -->`` sets the slide's speaker notes
- ``::: columns`` ... ``:::`` lays the enclosed blocks out side by side

Example:
    ```python
    from easypptx import Presentation

    pres = Presentation.from_markdown("deck.md")
    pres.save("deck.pptx")
    ```
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from easypptx.presentation import Presentation
    from easypptx.slide import Slide as EasySlide

_NOTES_RE = re.compile(r"<!--\s*notes?:\s*(.*?)\s*-->", re.DOTALL)
_IMAGE_RE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")
_BULLET_RE = re.compile(r"^(\s*)(?:[-*+]|\d+[.)])\s+(.*)$")

# Layout constants (percent of slide)
_TITLE_BAND = 14
_CONTENT_TOP = 18
_CONTENT_BOTTOM = 95
_GAP = 2


@dataclass
class _Block:
    """One renderable content block on a slide."""

    kind: str  # "bullets" | "text" | "image" | "table" | "code" | "columns"
    lines: list = field(default_factory=list)
    items: list = field(default_factory=list)  # bullets: (text, level)
    columns: list = field(default_factory=list)  # columns: list[_Block]

    @property
    def weight(self) -> float:
        """Relative share of vertical space this block wants."""
        if self.kind == "bullets":
            return max(2.0, len(self.items) * 0.9)
        if self.kind == "image":
            return 6.0
        if self.kind == "table":
            return max(2.0, len(self.lines) * 0.8)
        if self.kind == "code":
            return max(2.0, len(self.lines) * 0.6)
        if self.kind == "columns":
            return max((b.weight for b in self.columns), default=2.0)
        return 1.5  # text paragraph


@dataclass
class _SlideSpec:
    """Parsed content for one slide."""

    title: str | None = None
    blocks: list[_Block] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    is_title_slide: bool = False
    subtitle: str | None = None


def _parse_frontmatter(lines: list[str]) -> tuple[dict[str, str], int]:
    """Parse simple ``key: value`` frontmatter; returns (settings, lines consumed)."""
    if not lines or lines[0].strip() != "---":
        return {}, 0
    settings: dict[str, str] = {}
    for i, line in enumerate(lines[1:], start=1):
        if line.strip() == "---":
            return settings, i + 1
        if ":" in line:
            key, _, value = line.partition(":")
            settings[key.strip().lower()] = value.strip().strip("\"'")
    return {}, 0  # unterminated: treat as normal content


def _parse(text: str) -> tuple[dict[str, str], list[_SlideSpec]]:
    """Parse markdown into frontmatter settings and slide specs."""
    lines = text.splitlines()
    settings, consumed = _parse_frontmatter(lines)
    lines = lines[consumed:]

    slides: list[_SlideSpec] = []
    current: _SlideSpec | None = None
    bullets: list[tuple[str, int]] = []
    paragraph: list[str] = []
    in_code = False
    code_lines: list[str] = []
    table_lines: list[str] = []
    columns_mode = False
    column_blocks: list[_Block] = []

    def slide() -> _SlideSpec:
        nonlocal current
        if current is None:
            current = _SlideSpec()
            slides.append(current)
        return current

    def sink() -> list[_Block]:
        """Where finished blocks go: the current columns group or the slide."""
        return column_blocks if columns_mode else slide().blocks

    def flush_bullets() -> None:
        nonlocal bullets
        if bullets:
            sink().append(_Block("bullets", items=bullets))
            bullets = []

    def flush_paragraph() -> None:
        nonlocal paragraph
        if paragraph:
            sink().append(_Block("text", lines=[" ".join(paragraph)]))
            paragraph = []

    def flush_table() -> None:
        nonlocal table_lines
        if table_lines:
            sink().append(_Block("table", lines=table_lines))
            table_lines = []

    def flush_all() -> None:
        flush_bullets()
        flush_paragraph()
        flush_table()

    def end_columns() -> None:
        nonlocal columns_mode, column_blocks
        if columns_mode:
            flush_all()
            columns_mode = False
            if column_blocks:
                slide().blocks.append(_Block("columns", columns=column_blocks))
            column_blocks = []

    for raw in lines:
        line = raw.rstrip("\n")
        stripped = line.strip()

        # Fenced code blocks swallow everything until the closing fence
        if in_code:
            if stripped.startswith("```"):
                sink().append(_Block("code", lines=code_lines))
                code_lines = []
                in_code = False
            else:
                code_lines.append(line)
            continue
        if stripped.startswith("```"):
            flush_all()
            in_code = True
            continue

        # Speaker notes comments (attached after any heading creates its slide)
        notes_match = _NOTES_RE.search(line)
        note_text = None
        if notes_match:
            note_text = notes_match.group(1)
            line = _NOTES_RE.sub("", line)
            stripped = line.strip()

        # Column layout fences
        if stripped.startswith(":::"):
            marker = stripped.lstrip(":").strip().lower()
            if marker in ("columns", "cols"):
                flush_all()
                columns_mode = True
            else:
                end_columns()
            continue

        # Headings
        if stripped.startswith("#"):
            end_columns()
            flush_all()
            level = len(stripped) - len(stripped.lstrip("#"))
            heading = stripped[level:].strip()
            if level == 1 and not slides:
                slides.append(_SlideSpec(title=heading, is_title_slide=True))
                current = slides[-1]
            else:
                current = _SlideSpec(title=heading)
                slides.append(current)
            if note_text:
                current.notes.append(note_text)
            continue

        if note_text:
            slide().notes.append(note_text)
            if not stripped:
                continue

        # Horizontal rule = forced slide break
        if re.fullmatch(r"(-{3,}|\*{3,}|_{3,})", stripped):
            end_columns()
            flush_all()
            current = _SlideSpec()
            slides.append(current)
            continue

        # Blank line closes paragraph/bullet/table groups
        if not stripped:
            flush_all()
            continue

        # Images on their own line
        image_match = _IMAGE_RE.match(stripped)
        if image_match:
            flush_all()
            sink().append(_Block("image", lines=[image_match.group(2)]))
            continue

        # Tables
        if stripped.startswith("|") and stripped.endswith("|"):
            flush_bullets()
            flush_paragraph()
            table_lines.append(stripped)
            continue
        flush_table()

        # Bullets (2-space indent per nesting level)
        bullet_match = _BULLET_RE.match(line)
        if bullet_match:
            flush_paragraph()
            indent = len(bullet_match.group(1).replace("\t", "  "))
            bullets.append((bullet_match.group(2).strip(), min(4, indent // 2)))
            continue
        flush_bullets()

        # The first paragraph on a title slide becomes the subtitle
        spec = slide()
        if spec.is_title_slide and not spec.blocks and spec.subtitle is None and not paragraph:
            spec.subtitle = stripped
            continue

        paragraph.append(stripped)

    end_columns()
    flush_all()
    return settings, slides


def _parse_table(lines: list[str]) -> list[list[str]]:
    """Convert GFM table lines into a list of rows, dropping the separator."""
    rows = []
    for line in lines:
        cells = [c.strip() for c in line.strip("|").split("|")]
        if all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c):
            continue  # separator row
        rows.append(cells)
    return rows


def _strip_inline(text: str) -> str:
    """Drop basic inline markdown (bold/italic/code/link syntax)."""
    text = re.sub(r"\*\*(.+?)\*\*|__(.+?)__", lambda m: m.group(1) or m.group(2), text)
    text = re.sub(r"\*(.+?)\*|_(.+?)_", lambda m: m.group(1) or m.group(2), text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    return text


def _render_block(
    slide: EasySlide, block: _Block, x: float, y: float, width: float, height: float, base_dir: Path
) -> None:
    """Render one block into the given slide region (percent coordinates)."""
    if block.kind == "bullets":
        items: list[str | tuple[str, int]] = [(_strip_inline(t), lvl) for t, lvl in block.items]
        slide.add_bullets(items, x=x, y=y, width=width, height=height)
    elif block.kind == "text":
        slide.add_text(_strip_inline(block.lines[0]), x=x, y=y, width=width, height=height)
    elif block.kind == "image":
        path = Path(block.lines[0])
        if not path.is_absolute():
            path = base_dir / path
        from easypptx.image import Image

        Image(slide).add(image_path=path, x=x, y=y, width=width, height=height)
    elif block.kind == "table":
        rows = _parse_table(block.lines)
        if rows:
            slide.add_table(rows, x=x, y=y, width=width, height=height, has_header=True)
    elif block.kind == "code":
        shape = slide.add_text(
            "\n".join(block.lines) if block.lines else "",
            x=x,
            y=y,
            width=width,
            height=height,
            font_name="Consolas",
            font_size=12,
        )
        # Code blocks keep every line as its own paragraph
        from pptx.util import Pt

        frame = shape.text_frame
        frame.clear()
        for i, code_line in enumerate(block.lines or [""]):
            paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            paragraph.text = code_line
            paragraph.font.name = "Consolas"
            paragraph.font.size = Pt(12)
    elif block.kind == "columns":
        n = len(block.columns)
        if n == 0:
            return
        gap = _GAP
        col_width = (width - gap * (n - 1)) / n
        for i, column in enumerate(block.columns):
            _render_block(slide, column, x + i * (col_width + gap), y, col_width, height, base_dir)


def _render_slide(pres: Presentation, spec: _SlideSpec, base_dir: Path) -> EasySlide:
    """Render one parsed slide spec into the presentation."""
    if spec.is_title_slide:
        slide = pres.add_slide()
        slide.add_text(
            spec.title or "",
            x=5,
            y=30,
            width=90,
            height=18,
            font_size=40,
            font_bold=True,
            align="center",
            vertical="middle",
        )
        if spec.subtitle:
            slide.add_text(_strip_inline(spec.subtitle), x=10, y=52, width=80, height=10, font_size=20, align="center")
    else:
        slide = pres.add_slide(title=spec.title) if spec.title else pres.add_slide()
        top = _CONTENT_TOP if spec.title else 8
        available = _CONTENT_BOTTOM - top

        blocks = spec.blocks
        if blocks:
            total_weight = sum(b.weight for b in blocks)
            gaps = _GAP * (len(blocks) - 1)
            usable = available - gaps
            y = float(top)
            for block in blocks:
                block_height = usable * block.weight / total_weight
                _render_block(slide, block, x=5, y=y, width=90, height=block_height, base_dir=base_dir)
                y += block_height + _GAP

    if spec.notes:
        slide.notes = "\n".join(spec.notes)
    return slide


def from_markdown(
    source: str | Path,
    theme: Any = None,
    template_toml: str | None = None,
    base_dir: str | Path | None = None,
) -> Presentation:
    """Build a Presentation from markdown text or a markdown file.

    Args:
        source: Markdown content, or a path to a .md file
        theme: Theme name or Theme instance; overrides the frontmatter theme (default: None)
        template_toml: TOML template path; overrides the frontmatter template (default: None)
        base_dir: Directory for resolving relative image paths
            (default: the markdown file's directory, or the current directory)

    Returns:
        A Presentation with one slide per markdown section

    Raises:
        FileNotFoundError: If source is a path to a missing file
    """
    from easypptx.presentation import Presentation

    path: Path | None = None
    if isinstance(source, Path):
        path = source
    elif "\n" not in source and source.strip().endswith((".md", ".markdown")):
        path = Path(source)

    if path is not None:
        if not path.exists():
            raise FileNotFoundError(f"Markdown file not found: {path}")
        text = path.read_text(encoding="utf-8")
        resolved_base = Path(base_dir) if base_dir is not None else path.parent
    else:
        text = str(source)
        resolved_base = Path(base_dir) if base_dir is not None else Path.cwd()

    settings, slides = _parse(text)

    theme = theme if theme is not None else settings.get("theme")
    template_toml = template_toml if template_toml is not None else settings.get("template")
    aspect_ratio = settings.get("aspect_ratio", "16:9")

    pres = Presentation(aspect_ratio=aspect_ratio, theme=theme, template_toml=template_toml)
    for spec in slides:
        _render_slide(pres, spec, resolved_base)
    return pres
