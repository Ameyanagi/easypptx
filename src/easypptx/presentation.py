"""Core presentation module for EasyPPTX."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any, ClassVar, Literal, cast

from pptx import Presentation as PPTXPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.autoshape import Shape as PPTXShape
from pptx.util import Inches

from easypptx import common
from easypptx.grid import Grid
from easypptx.image import Image
from easypptx.positioning import (
    apply_content_padding,
    pct,
    resolve_padding,
    shift_band,
    to_percent,
)
from easypptx.pyplot import Pyplot
from easypptx.slide import Slide
from easypptx.styles import Theme, resolve_theme
from easypptx.table import Table
from easypptx.template import Template, TemplateManager
from easypptx.text import Text

if TYPE_CHECKING:
    import pandas as pd


class Presentation:
    """Main presentation class for creating and manipulating PowerPoint presentations.

    This class provides a simplified interface for working with PowerPoint presentations,
    making it easy to create, modify, and save PPTX files.

    For standard aspect ratios (16:9 and 4:3), reference templates are automatically used
    when no specific template is provided, ensuring consistent and attractive presentations.

    Attributes:
        pptx_presentation: The underlying python-pptx Presentation object

    Examples:
        ```python
        # Create a new presentation with default 16:9 aspect ratio
        # (automatically uses the reference_16x9.pptx template)
        presentation = Presentation()

        # Create a presentation with 4:3 aspect ratio
        # (automatically uses the reference_4x3.pptx template)
        presentation = Presentation(aspect_ratio="4:3")

        # Create a presentation with custom dimensions
        # (doesn't use reference templates)
        presentation = Presentation(width_inches=13.33, height_inches=7.5)

        # Create a presentation with a custom template
        # (overrides the reference templates)
        presentation = Presentation(template_path="my_template.pptx")

        # Open an existing presentation
        presentation = Presentation.open("example.pptx")

        # Add a new slide
        slide = presentation.add_slide()

        # Save the presentation
        presentation.save("output.pptx")
        ```
    """

    # Standard aspect ratios in width:height format
    ASPECT_RATIOS: ClassVar[dict[str, tuple[float, float]]] = {
        "16:9": (13.33, 7.5),  # Widescreen (default)
        "4:3": (10, 7.5),  # Standard
        "16:10": (13.33, 8.33),  # Widescreen alternative
        "A4": (11.69, 8.27),  # A4 paper size
        "LETTER": (11, 8.5),  # US Letter paper size
    }

    # Shared constants (single source of truth in easypptx.common)
    COLORS: ClassVar[dict[str, RGBColor]] = common.COLORS
    ALIGN: ClassVar[dict[str, PP_ALIGN]] = common.ALIGN
    VERTICAL: ClassVar[dict[str, MSO_ANCHOR]] = common.VERTICAL
    DEFAULT_FONT = common.DEFAULT_FONT

    def __init__(
        self,
        aspect_ratio: str | None = "16:9",
        width_inches: float | None = None,
        height_inches: float | None = None,
        template_path: str | None = None,
        reference_pptx: str | None = None,
        blank_layout_index: int | None = None,
        default_bg_color: str | tuple[int, int, int] | None = None,
        template_toml: str | None = None,
        theme: str | Theme | None = None,
    ) -> None:
        """Initialize a new empty presentation.

        Args:
            aspect_ratio: Predefined aspect ratio, one of "16:9" (default), "4:3", "16:10", "A4", "LETTER"
            width_inches: Custom width in inches (overrides aspect_ratio if specified)
            height_inches: Custom height in inches (overrides aspect_ratio if specified)
            template_path: Path to a reference PowerPoint template to use for styles (default: None)
            reference_pptx: Path to a custom reference PPTX file to use (default: None)
            blank_layout_index: Index of blank layout in the slide_layouts (default: None, auto-detected)
            default_bg_color: Default background color for slides as string name or RGB tuple (default: None)
            template_toml: Path to a TOML template file to use for slides (default: None)
            theme: Built-in theme name ("light", "dark", "corporate") or a
                Theme instance applied to every slide (default: None)

        Raises:
            ValueError: If an invalid aspect ratio is specified
            FileNotFoundError: If the template file doesn't exist
        """
        # Late chance to register df.pptx (pandas may have been imported
        # after easypptx)
        from easypptx.pandas_accessor import maybe_register

        maybe_register()

        self.theme = resolve_theme(theme)
        if default_bg_color is None and self.theme is not None:
            default_bg_color = self.theme.bg_color
        self.default_bg_color = default_bg_color

        # Initialize the Template object and TemplateManager
        self.template = Template()
        self.template_manager = TemplateManager()

        # Track which reference PPTX file we've loaded
        self._loaded_reference = None

        # Store default template name if a TOML template is provided
        self._default_template = None

        # Cache of Slide wrappers keyed by slide id, so repeated access
        # returns the same object (preserving user_data and identity)
        self._slide_cache: dict[int, Slide] = {}

        # Load TOML template if provided (raises FileNotFoundError/ValueError)
        if template_toml:
            self._default_template = self.template_manager.load(template_toml)

        if template_path:
            # Use an existing template
            if not Path(template_path).exists():
                raise FileNotFoundError(f"Template file not found: {template_path}")
            try:
                self.pptx_presentation = PPTXPresentation(template_path)
                self._loaded_reference = str(template_path)
            except Exception as e:
                raise ValueError(f"Invalid template file '{template_path}': {e}") from e
        elif reference_pptx:
            # Use a custom reference PPTX file
            reference_path = Path(reference_pptx)
            if not reference_path.exists():
                raise FileNotFoundError(f"Reference PPTX file not found: {reference_pptx}")

            try:
                self.pptx_presentation = PPTXPresentation(str(reference_path))
                self._loaded_reference = str(reference_path)
            except Exception as e:
                raise ValueError(f"Invalid reference PPTX file '{reference_pptx}': {e}") from e
        else:
            # Check if we should use a reference template based on aspect ratio
            reference_template = None

            # Only use reference templates for the specific aspect ratios we have templates for
            if aspect_ratio == "16:9" and width_inches is None and height_inches is None:
                # Use 16:9 reference template
                reference_template = Path(__file__).parent / "reference_16x9.pptx"
            elif aspect_ratio == "4:3" and width_inches is None and height_inches is None:
                # Use 4:3 reference template
                reference_template = Path(__file__).parent / "reference_4x3.pptx"

            if reference_template and reference_template.exists():
                # Use the appropriate reference template
                self.pptx_presentation = PPTXPresentation(str(reference_template))
                self._loaded_reference = str(reference_template)
            else:
                # Create a new presentation without a template
                self.pptx_presentation = PPTXPresentation()

            # Set slide dimensions based on inputs
            if width_inches is not None and height_inches is not None:
                # Use custom dimensions
                self._set_slide_dimensions(width_inches, height_inches)
            elif aspect_ratio is not None:
                # Use predefined aspect ratio
                if aspect_ratio not in self.ASPECT_RATIOS:
                    valid_ratios = ", ".join(self.ASPECT_RATIOS.keys())
                    raise ValueError(f"Invalid aspect ratio: {aspect_ratio}. Valid options are: {valid_ratios}")

                width, height = self.ASPECT_RATIOS[aspect_ratio]
                self._set_slide_dimensions(width, height)

        # Find and store the blank slide layout
        if blank_layout_index is not None:
            # Use the specified index
            if 0 <= blank_layout_index < len(self.pptx_presentation.slide_layouts):
                self.blank_layout = self.pptx_presentation.slide_layouts[blank_layout_index]
            else:
                # If index is out of range, fall back to a safe default
                self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[0]
        else:
            # Auto-detect the blank layout (typically index 6, but can vary)
            self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[6]

    @property
    def _slide_width_emu(self) -> int:
        """Slide width in EMUs (falls back to 16:9 default)."""
        width = self.pptx_presentation.slide_width
        return int(width) if width is not None else 12192768

    @property
    def _slide_height_emu(self) -> int:
        """Slide height in EMUs (falls back to 16:9 default)."""
        height = self.pptx_presentation.slide_height
        return int(height) if height is not None else 6858000

    def _wrap_slide(self, pptx_slide: Any) -> Slide:
        """Return the cached Slide wrapper for a python-pptx slide, creating it if needed."""
        key = pptx_slide.slide_id
        slide = self._slide_cache.get(key)
        if slide is None:
            slide = Slide(
                pptx_slide,
                slide_width=self._slide_width_emu,
                slide_height=self._slide_height_emu,
            )
            if self.theme is not None:
                slide.apply_template_defaults(self.theme.to_template())
            self._slide_cache[key] = slide
        return slide

    def _find_blank_layout(self) -> Any:
        """Find the blank layout in the presentation.

        Attempts to find the layout with the fewest placeholders, which is typically the blank layout.

        Returns:
            The slide layout that appears to be blank, or None if no suitable layout is found
        """
        # First, check if there's a layout named "Blank" or similar
        for layout in self.pptx_presentation.slide_layouts:
            if hasattr(layout, "name") and "blank" in layout.name.lower():
                return layout

        # Find the layout with the fewest placeholders (likely to be blank)
        blank_layout = None
        min_placeholders = float("inf")

        for layout in self.pptx_presentation.slide_layouts:
            placeholder_count = len(layout.placeholders)
            if placeholder_count < min_placeholders:
                min_placeholders = placeholder_count
                blank_layout = layout

        return blank_layout if min_placeholders < 3 else None

    def _set_slide_dimensions(self, width_inches: float, height_inches: float) -> None:
        """Set the slide dimensions.

        Args:
            width_inches: Width in inches
            height_inches: Height in inches
        """
        self.pptx_presentation.slide_width = Inches(width_inches)
        self.pptx_presentation.slide_height = Inches(height_inches)

    def deck(self) -> Any:
        """Start a fluent Deck builder on this presentation.

        Returns:
            A Deck whose chained calls render onto this presentation at save time

        Examples:
            ```python
            pres = Presentation(theme="dark")
            pres.deck().slide("Agenda").bullets(["One", "Two"]).save("out.pptx")
            ```
        """
        from easypptx.deck import Deck

        return Deck(presentation=self)

    @classmethod
    def from_markdown(
        cls,
        source: str | Path,
        theme: str | Theme | None = None,
        template_toml: str | None = None,
        base_dir: str | Path | None = None,
    ) -> Presentation:
        """Build a presentation from markdown text or a markdown file.

        See :func:`easypptx.markdown.from_markdown` for the supported syntax.

        Args:
            source: Markdown content, or a path to a .md file
            theme: Theme name or Theme instance; overrides the frontmatter theme (default: None)
            template_toml: TOML template path; overrides the frontmatter template (default: None)
            base_dir: Directory for resolving relative image paths (default: the file's directory)

        Returns:
            A Presentation with one slide per markdown section

        Examples:
            ```python
            pres = Presentation.from_markdown("deck.md")
            pres.save("deck.pptx")
            ```
        """
        from easypptx.markdown import from_markdown

        return from_markdown(source, theme=theme, template_toml=template_toml, base_dir=base_dir)

    @classmethod
    def open(cls, file_path: str | Path, blank_layout_index: int | None = None) -> Presentation:
        """Open an existing PowerPoint presentation.

        Args:
            file_path: Path to the PowerPoint file to open
            blank_layout_index: Index of blank layout in the slide_layouts (default: None, auto-detected)

        Returns:
            A new Presentation object with the loaded presentation

        Raises:
            FileNotFoundError: If the specified file doesn't exist
            ValueError: If the file is not a valid PowerPoint file
        """
        file_path_obj = Path(file_path)

        if not file_path_obj.exists():
            raise FileNotFoundError(f"Presentation file not found: {file_path}")

        try:
            pptx_presentation = PPTXPresentation(str(file_path_obj))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint file: {e}") from e

        presentation = cls(width_inches=None, height_inches=None, blank_layout_index=blank_layout_index)
        presentation.pptx_presentation = pptx_presentation
        presentation._loaded_reference = str(file_path_obj)

        # Find and store the blank slide layout
        if blank_layout_index is not None:
            # Use the specified index
            if 0 <= blank_layout_index < len(presentation.pptx_presentation.slide_layouts):
                presentation.blank_layout = presentation.pptx_presentation.slide_layouts[blank_layout_index]
            else:
                # If index is out of range, fall back to a safe default
                presentation.blank_layout = (
                    presentation._find_blank_layout() or presentation.pptx_presentation.slide_layouts[0]
                )
        else:
            # Auto-detect the blank layout (typically index 6, but can vary)
            presentation.blank_layout = (
                presentation._find_blank_layout() or presentation.pptx_presentation.slide_layouts[6]
            )

        return presentation

    def add_slide(
        self,
        layout_index: int | None = None,
        bg_color: str | tuple[int, int, int] | None = None,
        title: str | None = None,
        template_toml: str | Literal[False] | None = None,
        title_padding: str | float | None = None,
        title_x_padding: str | float | None = "5%",
        title_y_padding: str | float | None = "5%",
        title_width: str | float | None = "90%",
        title_height: str | float | None = "10%",
        title_font_size: int | None = None,
        title_align: str | None = None,
        title_color: str | tuple[int, int, int] | None = None,
    ) -> Slide:
        """Add a new slide to the presentation.

        Args:
            layout_index: Index of the slide layout to use (default: None uses blank layout)
            bg_color: Background color for this slide, overrides default (default: None)
            title: Optional title text for the slide (default: None)
            template_toml: Path to a TOML template file to use for this slide.
                Pass False to opt out of the presentation's default template
                for this slide (default: None, uses the default template if set)
            title_padding: Padding around the title (applies to both x and y if specified) (default: None)
            title_x_padding: Horizontal padding for title (default: "5%")
            title_y_padding: Vertical padding for title (default: "5%")
            title_width: Width of the title area (default: "90%")
            title_height: Height of the title area (default: "10%")
            title_font_size: Font size for the title (default: theme title size or 32)
            title_align: Text alignment for the title (default: theme title align or "center")
            title_color: Title color as a name from COLORS or an RGB tuple
                (default: theme title color or None)

        Returns:
            A new Slide object

        Raises:
            FileNotFoundError: If template_toml points to a missing file
            ValueError: If the template file is invalid
        """
        # Priority for templates:
        # 1. If template_toml is explicitly provided for this slide, use it
        # 2. If template_toml is False, use no template for this slide
        # 3. Otherwise, if there's a default template from __init__, use that
        # 4. Finally, fall back to a standard slide

        template_name = None

        # If a template is directly provided for this slide
        if template_toml:
            template_name = self.template_manager.load(template_toml)

        # Use the default template unless templates are opted out for this slide
        if template_name is None and template_toml is not False and self._default_template is not None:
            template_name = self._default_template

        # If we have a template (either slide-specific or default), use it
        if template_name is not None:
            # Use the template
            slide = self.add_slide_from_template(template_name)

            # Apply title if provided
            if title is not None:
                # Find the first text frame to use as title
                title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
                if title_shapes and len(title_shapes) > 0:
                    title_shape = title_shapes[0]
                    title_shape.text_frame.text = title

            # Apply custom background color if specified
            if bg_color is not None:
                slide.set_background_color(bg_color)

            return slide
        else:
            # Use blank layout by default, or specified layout if provided
            slide_layout = (
                self.blank_layout if layout_index is None else self.pptx_presentation.slide_layouts[layout_index]
            )

            pptx_slide = self.pptx_presentation.slides.add_slide(slide_layout)
            slide = self._wrap_slide(pptx_slide)

            # Apply background color if specified for this slide or as default
            color_to_use = bg_color if bg_color is not None else self.default_bg_color
            if color_to_use is not None:
                slide.set_background_color(color_to_use)

            # Add title if provided
            if title is not None:
                # Resolve title styling: explicit > theme > library default
                title_style = self.theme.title if self.theme is not None else None
                if title_font_size is None:
                    title_font_size = (title_style.font_size if title_style else None) or 32
                if title_align is None:
                    title_align = (title_style.align if title_style else None) or "center"
                if title_color is None and title_style is not None:
                    title_color = title_style.color
                title_font_name = title_style.font_name if title_style else None
                title_font_bold = title_style.font_bold if title_style and title_style.font_bold is not None else True

                # If general title_padding is provided, it overrides individual x and y padding
                x_padding = title_padding if title_padding is not None else title_x_padding
                y_padding = title_padding if title_padding is not None else title_y_padding

                # We don't have add_title method, so use add_text with appropriate positioning
                slide.add_text(
                    text=title,
                    x=x_padding,
                    y=y_padding,
                    width=title_width,
                    height=title_height,
                    font_size=title_font_size,
                    font_bold=title_font_bold,
                    font_name=title_font_name,
                    align=title_align,
                    color=title_color,
                )

                # A short accent bar under the title anchors the visual hierarchy
                if self.theme is not None and self.theme.title_accent and self.theme.accent_color is not None:
                    bar_x = to_percent(x_padding if x_padding is not None else "5%", self._slide_width_emu)
                    bar_y = to_percent(
                        y_padding if y_padding is not None else "5%", self._slide_height_emu
                    ) + to_percent(title_height if title_height is not None else "10%", self._slide_height_emu)
                    bar = slide.add_shape(
                        x=pct(bar_x),
                        y=pct(bar_y),
                        width=7,
                        height=0.9,
                        fill_color=self.theme.accent_color,
                    )
                    bar.line.fill.background()

            return slide

    @property
    def slides(self) -> list[Slide]:
        """Get a list of all slides in the presentation.

        Returns:
            List of Slide objects. Repeated access returns the same wrapper
            objects, so state like user_data is preserved.
        """
        return [self._wrap_slide(slide) for slide in self.pptx_presentation.slides]

    def save(self, file_path: str | Path) -> None:
        """Save the presentation to a file.

        Args:
            file_path: Path where the presentation should be saved
        """
        self.pptx_presentation.save(str(file_path))

    def add_slide_from_template(self, template_data: str | dict) -> Slide:
        """Add a slide using a template preset.

        Args:
            template_data: Template preset name or custom template dictionary

        Returns:
            A new Slide object configured according to the template
        """
        # Get the template data - either from a preset or use the provided dictionary
        template_name = None
        if isinstance(template_data, str):
            template_name = template_data
            try:
                # First try to get from TemplateManager (includes built-in and registered templates)
                preset = self.template_manager.get(template_data)
            except ValueError:
                # Fall back to the legacy approach if not found
                preset = self.template.get_preset(template_data)
        else:
            preset = template_data

        # Check if there's a reference PPTX file specified for this template
        reference_pptx = None
        blank_layout_index = None
        if template_name is not None:
            # First check in TemplateManager for loaded templates
            reference_pptx = self.template_manager.get_reference_pptx(template_name)
            blank_layout_index = self.template_manager.get_blank_layout_index(template_name)

            # If not found, check built-in presets
            if reference_pptx is None:
                reference_pptx = self.template.get_reference_pptx(template_name)
                blank_layout_index = self.template.get_blank_layout_index(template_name)

            # If a reference PPTX is specified and we haven't loaded it yet, switch to it.
            # This replaces the underlying presentation, so it is only allowed
            # before any slides have been added. Compare resolved paths so the
            # same file referenced two ways doesn't trigger a needless swap.
            already_loaded = (
                reference_pptx is not None
                and self._loaded_reference is not None
                and Path(self._loaded_reference).resolve() == Path(reference_pptx).resolve()
            )
            if reference_pptx is not None and not already_loaded:
                if len(self.pptx_presentation.slides) > 0:
                    raise ValueError(
                        f"Template requires reference PPTX '{reference_pptx}', but the presentation "
                        "already contains slides. Pass the reference PPTX (or template_toml) to "
                        "Presentation() at construction instead."
                    )

                # Save current properties we want to preserve
                current_width = self.pptx_presentation.slide_width
                current_height = self.pptx_presentation.slide_height

                # Load the reference PPTX
                try:
                    self.pptx_presentation = PPTXPresentation(reference_pptx)
                except Exception as e:
                    raise ValueError(f"Invalid reference PPTX file '{reference_pptx}': {e}") from e
                self._loaded_reference = reference_pptx
                self._slide_cache.clear()

                # Restore dimensions
                if current_width is not None:
                    self.pptx_presentation.slide_width = current_width
                if current_height is not None:
                    self.pptx_presentation.slide_height = current_height

                # Update blank layout
                if blank_layout_index is not None and 0 <= blank_layout_index < len(
                    self.pptx_presentation.slide_layouts
                ):
                    self.blank_layout = self.pptx_presentation.slide_layouts[blank_layout_index]
                else:
                    self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[6]

        # Get background color if specified
        bg_color = preset.get("bg_color", None)

        # Convert list colors to tuples for compatibility with slide.set_background_color
        if isinstance(bg_color, list) and len(bg_color) == 3:
            bg_color = tuple(bg_color)

        # Create a new slide using the blank layout directly (avoiding recursion)
        # Use the blank layout or specified layout if provided
        slide_layout = self.blank_layout
        pptx_slide = self.pptx_presentation.slides.add_slide(slide_layout)
        slide = self._wrap_slide(pptx_slide)

        # Apply background color if specified for this slide
        if bg_color is not None:
            slide.set_background_color(bg_color)

        # Add title if specified in template
        if "title" in preset:
            title_data = preset["title"]
            position = title_data.get("position", {"x": "5%", "y": "5%", "width": "90%", "height": "10%"})

            # Extract font information
            font_data = title_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 32)
            font_bold = font_data.get("bold", True)

            # Extract alignment information
            align = title_data.get("align", "center")
            vertical = title_data.get("vertical", "middle")
            color = title_data.get("color", "black")

            # Convert list colors to tuples
            if isinstance(color, list) and len(color) == 3:
                color = tuple(color)

            # Add the title text
            Text.add(
                slide=slide,
                text=title_data.get("text", "Title"),
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        # Add subtitle if specified in template
        if "subtitle" in preset:
            subtitle_data = preset["subtitle"]
            position = subtitle_data.get("position", {"x": "20%", "y": "60%", "width": "60%", "height": "20%"})

            # Extract font information
            font_data = subtitle_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 24)
            font_bold = font_data.get("bold", False)

            # Extract alignment information
            align = subtitle_data.get("align", "center")
            vertical = subtitle_data.get("vertical", "middle")
            color = subtitle_data.get("color", "black")

            # Convert list colors to tuples
            if isinstance(color, list) and len(color) == 3:
                color = tuple(color)

            # Add the subtitle text
            Text.add(
                slide=slide,
                text=subtitle_data.get("text", "Subtitle"),
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        # Add decorative bar if specified
        if "bar" in preset:
            bar_data = preset["bar"]
            position = bar_data.get("position", {"x": "0%", "y": "10%", "width": "100%", "height": "2%"})

            # Create rectangle shape for the bar
            shape = slide.add_shape(
                shape_type=MSO_SHAPE.RECTANGLE,
                x=position.get("x", "0%"),
                y=position.get("y", "10%"),
                width=position.get("width", "100%"),
                height=position.get("height", "2%"),
            )

            # Apply gradient if specified
            if "gradient" in bar_data:
                gradient = bar_data["gradient"]
                start_color = gradient.get("start_color")
                end_color = gradient.get("end_color")
                angle = gradient.get("angle", 0)

                # Convert list colors to tuples
                if isinstance(start_color, list) and len(start_color) == 3:
                    start_color = tuple(start_color)
                if isinstance(end_color, list) and len(end_color) == 3:
                    end_color = tuple(end_color)

                # Convert tuple colors to RGBColor objects
                if isinstance(start_color, tuple) and len(start_color) == 3:
                    start_color = RGBColor(*start_color)
                if isinstance(end_color, tuple) and len(end_color) == 3:
                    end_color = RGBColor(*end_color)

                fill = shape.fill
                fill.gradient()
                fill.gradient_stops[0].color.rgb = start_color
                fill.gradient_stops[1].color.rgb = end_color
                fill.gradient_angle = angle

        # Add content or other text elements as specified in template
        # (These will be filled in by the specific methods like add_content_slide, add_image_slide, etc.)

        # Store styling information in the slide's user data
        # This will be used by other methods like add_image_slide, add_table_slide, etc.
        slide.user_data = {
            "template_preset": preset,
            "image_style": self.template.get_image_style(preset) if "image_style" in preset else None,
            "table_style": self.template.get_table_style(preset) if "table_style" in preset else None,
            "chart_style": self.template.get_chart_style(preset) if "chart_style" in preset else None,
        }

        return slide

    def add_title_slide(self, title: str, subtitle: str | None = None) -> Slide:
        """Add a title slide with title and optional subtitle.

        Args:
            title: Text for the title
            subtitle: Text for the subtitle (default: None)

        Returns:
            A new Slide object configured as a title slide
        """
        # Get the title slide preset
        preset = self.template.get_preset("title_slide")

        # Create a new slide using the preset
        slide = self.add_slide_from_template("title_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Add or update the subtitle text if provided
        if subtitle and len(title_shapes) > 1:
            subtitle_shape = title_shapes[1]
            subtitle_shape.text_frame.text = subtitle
        elif subtitle:
            # Get subtitle data from preset
            subtitle_data = preset.get("subtitle", {})
            position = subtitle_data.get("position", {"x": "20%", "y": "60%", "width": "60%", "height": "20%"})

            # Extract font information
            font_data = subtitle_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 24)
            font_bold = font_data.get("bold", False)

            # Extract alignment information
            align = subtitle_data.get("align", "center")
            vertical = subtitle_data.get("vertical", "middle")
            color = subtitle_data.get("color", "black")

            # Convert list colors to tuples
            if isinstance(color, list) and len(color) == 3:
                color = tuple(color)

            # Add the subtitle text
            Text.add(
                slide=slide,
                text=subtitle,
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        return slide

    def add_content_slide(self, title: str, use_bar: bool = True) -> Slide:
        """Add a content slide with title and optional horizontal bar.

        Args:
            title: Text for the title
            use_bar: Whether to include decorative bar (default: True)

        Returns:
            A new Slide object configured as a content slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("content_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # If bar is not wanted, remove it
        if not use_bar:
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE.RECTANGLE and not shape.has_text_frame:
                    shape._element.getparent().remove(shape._element)

        return slide

    def add_section_slide(self, title: str, bg_color: str = "blue") -> Slide:
        """Add a section slide with a full-screen title on a colored background.

        Args:
            title: Text for the section title
            bg_color: Background color (default: "blue")

        Returns:
            A new Slide object configured as a section slide
        """
        # Create custom preset for section slide with specified background color
        preset = self.template.get_preset("section_slide")
        preset["bg_color"] = bg_color

        # Create a new slide using the preset
        slide = self.add_slide_from_template(preset)

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        return slide

    def add_image_gen_slide(
        self,
        image_path: str,
        title: str | None = None,
        subtitle: str | None = None,
        label: str | None = None,
        x: float | str = "10%",
        y: float | str = "20%",
        width: float | str = "80%",
        height: float | str = "70%",
        title_height: float | str = "10%",
        subtitle_height: float | str = "5%",
        bg_color: str | tuple[int, int, int] | None = None,
        title_font_size: int = 24,
        subtitle_font_size: int = 18,
        label_font_size: int = 14,
        border: bool = False,
        border_color: str = "black",
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
        title_padding: str | float | None = None,
        title_x_padding: str | float | None = None,
        title_y_padding: str | float | None = None,
        subtitle_padding: str | float | None = None,
        subtitle_x_padding: str | float | None = None,
        subtitle_y_padding: str | float | None = None,
        content_padding: str | float | None = None,
        content_x_padding: str | float | None = None,
        content_y_padding: str | float | None = None,
        label_padding: str | float | None = None,
        label_x_padding: str | float | None = None,
        label_y_padding: str | float | None = "1%",
        title_align: str = "center",
        subtitle_align: str = "center",
        label_align: str = "center",
    ) -> tuple[Slide, PPTXShape]:
        """Add a slide with an image and optional title, subtitle, and label.

        This method provides a more flexible alternative to add_image_slide with
        similar parameters to add_grid_slide and add_pyplot_slide.

        Args:
            image_path: Path to the image file
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            label: Optional caption for the image, displayed below (default: None)
            x: X position of the image as percentage or absolute value (default: "10%")
            y: Y position of the image as percentage or absolute value (default: "20%")
            width: Width of the image as percentage or absolute value (default: "80%")
            height: Height of the image as percentage or absolute value (default: "70%")
            title_height: Height of the title area (default: "10%")
            subtitle_height: Height of the subtitle area (default: "5%")
            bg_color: Background color for the slide (default: None)
            title_font_size: Font size for the title (default: 24)
            subtitle_font_size: Font size for the subtitle (default: 18)
            label_font_size: Font size for the caption (default: 14)
            border: Whether to add a border around the image (default: False)
            border_color: Color for the border (default: "black")
            shadow: Whether to add a shadow effect to the image (default: False)
            maintain_aspect_ratio: Whether to maintain the image's aspect ratio (default: True)
            title_padding: Padding around the title, applies to both x and y (default: None)
            title_x_padding: Horizontal padding for title, overridden by title_padding if provided (default: None)
            title_y_padding: Vertical padding for title, overridden by title_padding if provided (default: None)
            subtitle_padding: Padding around the subtitle, applies to both x and y (default: None)
            subtitle_x_padding: Horizontal padding for subtitle, overridden by subtitle_padding if provided (default: None)
            subtitle_y_padding: Vertical padding for subtitle, overridden by subtitle_padding if provided (default: None)
            content_padding: Padding around the content area, applies to both x and y (default: None)
            content_x_padding: Horizontal padding for content, overridden by content_padding if provided (default: None)
            content_y_padding: Vertical padding for content, overridden by content_padding if provided (default: None)
            label_padding: Padding around the label, applies to both x and y (default: None)
            label_x_padding: Horizontal padding for label, overridden by label_padding if provided (default: None)
            label_y_padding: Vertical padding between content and label (default: "1%")
            title_align: Text alignment for the title (default: "center")
            subtitle_align: Text alignment for the subtitle (default: "center")
            label_align: Text alignment for the label (default: "center")

        Returns:
            A tuple containing (Slide, PPTXShape) where PPTXShape is the image shape

        Example:
            ```python
            # Add an image with title and subtitle
            slide, image = pres.add_image_gen_slide(
                image_path="path/to/image.jpg",
                title="Product Showcase",
                subtitle="Latest Design",
                label="Figure 1: Product Prototype",
                maintain_aspect_ratio=True,
                title_padding="5%",
                content_padding="2%"
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)
        slide_height_emu = self._slide_height_emu

        # Calculate positions and dimensions
        adjusted_y = y
        adjusted_height = height

        # Determine content x-padding
        content_x_pad = resolve_padding(content_padding, content_x_padding)
        content_x_val = content_x_pad if content_x_pad is not None else x

        # Add title if provided
        if title:
            title_x = resolve_padding(title_padding, title_x_padding)
            title_y = resolve_padding(title_padding, title_y_padding)

            slide.add_text(
                text=title,
                x=title_x if title_x is not None else "0%",
                y=title_y if title_y is not None else "0%",
                width="100%",
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Reserve the title band above the content area
            adjusted_y, adjusted_height = shift_band(y, height, title_height, slide_height_emu)

        # Add subtitle if provided
        if subtitle:
            subtitle_x = resolve_padding(subtitle_padding, subtitle_x_padding)
            subtitle_y = resolve_padding(subtitle_padding, subtitle_y_padding)

            slide.add_text(
                text=subtitle,
                x=subtitle_x if subtitle_x is not None else "0%",
                y=subtitle_y if subtitle_y is not None else adjusted_y,
                width="100%",
                height=subtitle_height,
                font_size=subtitle_font_size,
                align=subtitle_align,
                vertical="middle",
            )

            # Reserve the subtitle band above the content area
            adjusted_y, adjusted_height = shift_band(adjusted_y, adjusted_height, subtitle_height, slide_height_emu)

        # Apply content y padding if specified
        image_y, adjusted_height = apply_content_padding(
            adjusted_y, adjusted_height, content_padding, content_y_padding, slide_height_emu
        )

        # Add the image to the slide
        img = Image(slide)
        image_shape = img.add(
            image_path=image_path,
            x=content_x_val,
            y=image_y,
            width=width,
            height=adjusted_height,
            maintain_aspect_ratio=maintain_aspect_ratio,
        )

        # Apply styling to the image
        if border:
            image_shape.line.color.rgb = self.COLORS.get(border_color, self.COLORS["black"])
            image_shape.line.width = 1  # 1 point width for border

        # Apply shadow if specified
        if shadow:
            common.apply_shadow(image_shape)

        # Add label if specified
        if label:
            label_x = resolve_padding(label_padding, label_x_padding)
            label_y_pad = resolve_padding(label_padding, label_y_padding)

            # Place the label just below the image
            label_y = pct(
                to_percent(image_y, slide_height_emu)
                + to_percent(adjusted_height, slide_height_emu)
                + (to_percent(label_y_pad, slide_height_emu) if label_y_pad is not None else 1.0)
            )

            slide.add_text(
                text=label,
                x=label_x if label_x is not None else "0%",
                y=label_y,
                width="100%",
                height="5%",
                font_size=label_font_size,
                align=label_align,
                vertical="top",
            )

        return slide, image_shape

    def add_comparison_slide(self, title: str, content_texts: list[str]) -> Slide:
        """Add a slide with title and two or more content areas for comparison.

        Args:
            title: Text for the title
            content_texts: List of texts for comparison areas (typically 2)

        Returns:
            A new Slide object configured as a comparison slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("comparison_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("comparison_slide")

        # Add content texts
        if len(content_texts) >= 1:
            left_position = preset.get("left_content", {}).get(
                "position", {"x": "5%", "y": "20%", "width": "42%", "height": "75%"}
            )
            Text.add(
                slide=slide,
                text=content_texts[0],
                position=left_position,
                font_name=self.DEFAULT_FONT,
                font_size=16,
                font_bold=False,
                align="left",
                vertical_align="top",
                color="black",
            )

        if len(content_texts) >= 2:
            right_position = preset.get("right_content", {}).get(
                "position", {"x": "53%", "y": "20%", "width": "42%", "height": "75%"}
            )
            Text.add(
                slide=slide,
                text=content_texts[1],
                position=right_position,
                font_name=self.DEFAULT_FONT,
                font_size=16,
                font_bold=False,
                align="left",
                vertical_align="top",
                color="black",
            )

        return slide

    def add_table_slide(
        self,
        title: str,
        data: list[list[Any]] | pd.DataFrame,
        has_header: bool = True,
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a table.

        Args:
            title: Text for the title
            data: Table data as a list of lists or pandas DataFrame
            has_header: Whether the first row is a header (default: True)
            custom_style: Dictionary of style options for the table (default: None)

        Returns:
            A new Slide object configured as a table slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("table_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("table_slide")
        table_position = preset.get("table_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        table_style = self.template.get_table_style(preset)
        if custom_style:
            # Handle nested dictionaries
            if "first_row" in custom_style and "first_row" in table_style:
                table_style["first_row"].update(custom_style.get("first_row", {}))

            # Update all other keys
            for key, value in custom_style.items():
                if key != "first_row":
                    table_style[key] = value

        # Add the table with styling
        table = Table(slide)

        # Extract position values
        x = table_position.get("x", "10%")
        y = table_position.get("y", "20%")
        width = table_position.get("width", "80%")
        height = table_position.get("height", "60%")

        # Convert pandas DataFrame to list if needed
        if common.is_dataframe(data):
            df = cast("pd.DataFrame", data)
            table_data: list[list[Any]] = [list(df.columns), *df.values.tolist()]
        else:
            table_data = cast("list[list[Any]]", data)

        # Add the table
        table.add(
            data=table_data,
            x=x,
            y=y,
            width=width,
            height=height,
            first_row_header=has_header,
            style=1 if table_style else None,
        )

        return slide

    def add_chart_slide(
        self,
        title: str,
        data: list[list[Any]] | pd.DataFrame,
        chart_type: str | None = None,
        category_column: str | None = None,
        value_columns: str | list[str] | None = None,
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a chart.

        Args:
            title: Text for the title
            data: Chart data as a list of lists or pandas DataFrame
            chart_type: Type of chart (default: None uses preset's chart_type)
            category_column: Name or index of the column to use as categories (default: None)
            value_columns: Names or indices of columns to use as values (default: None)
            custom_style: Dictionary of style options for the chart (default: None)

        Returns:
            A new Slide object configured as a chart slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("chart_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("chart_slide")
        chart_position = preset.get("chart_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        chart_style = self.template.get_chart_style(preset)
        if custom_style:
            chart_style.update(custom_style)

        # Use chart_type from parameters or from style
        if chart_type:
            chart_style["chart_type"] = chart_type

        # Import Chart class here to avoid circular import
        from easypptx.chart import Chart, extract_chart_series

        # Create a Chart object and add the chart
        chart_obj = Chart(slide)

        # Extract categories and every requested value series
        categories, series = extract_chart_series(data, category_column, value_columns)

        # Extract position values
        x = chart_position.get("x", "10%")
        y = chart_position.get("y", "20%")
        width = chart_position.get("width", "80%")
        height = chart_position.get("height", "70%")

        chart = chart_obj.add(
            chart_type=chart_style.get("chart_type", "column"),
            categories=categories,
            series=series,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            has_legend=chart_style.get("has_legend", True),
        )

        # Apply additional styling if applicable
        chart_any: Any = chart
        if hasattr(chart, "format") and chart_style.get("has_border", True):
            chart_any.format.line.color.rgb = self.COLORS.get(
                chart_style.get("border_color", "black"), self.COLORS["black"]
            )

        # Apply custom palette if specified
        if "palette" in chart_style and hasattr(chart, "series"):
            for i, series in enumerate(chart.series):
                if i < len(chart_style["palette"]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = chart_style["palette"][i]

        return slide

    # Direct object API methods for adding content to slides

    def add_grid(
        self,
        slide: Slide,
        x: float | str = "0%",
        y: float | str = "0%",
        width: float | str = "100%",
        height: float | str = "100%",
        rows: int = 1,
        cols: int = 1,
        padding: float = 5.0,
    ) -> Grid:
        """Add a grid layout to a slide.

        Args:
            slide: The slide to add the grid to
            x: X position in inches or percentage (default: "0%")
            y: Y position in inches or percentage (default: "0%")
            width: Width in inches or percentage (default: "100%")
            height: Height in inches or percentage (default: "100%")
            rows: Number of rows in the grid (default: 1)
            cols: Number of columns in the grid (default: 1)
            padding: Padding between cells as percentage of cell size (default: 5.0)

        Returns:
            The created Grid object

        Example:
            ```python
            slide = pres.add_slide()
            grid = pres.add_grid(slide, rows=2, cols=2)

            grid.add_to_cell(
                row=0,
                col=0,
                content_func=slide.add_text,
                text="Top Left",
                font_size=24,
                align="center",
                vertical="middle",
            )
            ```
        """
        from easypptx.grid import Grid

        # Create the grid
        grid = Grid(
            parent=slide,
            x=x,
            y=y,
            width=width,
            height=height,
            rows=rows,
            cols=cols,
            padding=padding,
        )

        return grid

    def add_simple_grid_slide(
        self,
        rows: int = 1,
        cols: int = 1,
        title: str | None = None,
        title_height: float | str = "10%",
        padding: float = 5.0,
        bg_color: str | tuple[int, int, int] | None = None,
    ) -> tuple[Slide, Grid]:
        """Add a slide with a simple grid layout.

        The grid will fill the entire slide, except for the title area if a title is provided.
        This is a simplified version of add_grid_slide for backward compatibility.

        Args:
            rows: Number of rows in the grid (default: 1)
            cols: Number of columns in the grid (default: 1)
            title: Optional title for the slide (default: None)
            title_height: Height of the title area (default: "10%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            bg_color: Background color for the slide (default: None)

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            slide, grid = pres.add_simple_grid_slide(rows=2, cols=3, title="My Grid Slide")

            grid.add_to_cell(
                row=0,
                col=0,
                content_func=slide.add_text,
                text="Top Left",
                font_size=24,
                align="center",
                vertical="middle",
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # Add title if specified
        if title:
            slide.add_text(
                text=title,
                x="0%",
                y="0%",
                width="100%",
                height=title_height,
                font_size=24,
                font_bold=True,
                align="center",
                vertical="middle",
            )

            # Adjust grid position and height to account for title
            grid_y, grid_height = shift_band("0%", "100%", title_height, self._slide_height_emu)

            # Create the grid
            grid = self.add_grid(
                slide=slide,
                x="0%",
                y=grid_y,
                width="100%",
                height=grid_height,
                rows=rows,
                cols=cols,
                padding=padding,
            )
        else:
            # Create the grid with full slide dimensions
            grid = self.add_grid(
                slide=slide,
                x="0%",
                y="0%",
                width="100%",
                height="100%",
                rows=rows,
                cols=cols,
                padding=padding,
            )

        return slide, grid

    def add_autogrid(
        self,
        slide: Slide,
        content_funcs: list | None = None,
        rows: int | None = None,
        cols: int | None = None,
        x: float | str = "0%",
        y: float | str = "0%",
        width: float | str = "100%",
        height: float | str = "100%",
        padding: float = 5.0,
        title: str | None = None,
        title_height: float | str = "10%",
        title_align: str = "center",
        column_major: bool = True,  # Use column-major order by default
    ) -> Grid:
        """Add an autogrid layout to a slide.

        This method automatically places the provided content functions into a grid.
        If content_funcs is None, an empty grid is created that can be populated later.

        Args:
            slide: The slide to add the autogrid to
            content_funcs: List of functions that add content to the slide, or None for an empty grid
            rows: Number of rows (default: None, calculated automatically)
            cols: Number of columns (default: None, calculated automatically)
            x: X position in inches or percentage (default: "0%")
            y: Y position in inches or percentage (default: "0%")
            width: Width in inches or percentage (default: "100%")
            height: Height in inches or percentage (default: "100%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            title: Optional title for the grid (default: None)
            title_height: Height of the title area (default: "10%")
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            column_major: Whether to fill cells in column-major order (default: True)
                         When True, fills cells down columns first, resulting in a visual layout
                         that matches the specified rows and columns when content is added sequentially.
                         When False, fills cells across rows first.

        Returns:
            The created Grid object

        Example:
            ```python
            slide = pres.add_slide()

            # Option 1: With content functions
            def create_text1():
                return slide.add_text("Text 1")

            def create_text2():
                return slide.add_text("Text 2")

            content_funcs = [create_text1, create_text2]
            pres.add_autogrid(slide, content_funcs, title="Auto Grid Example")

            # Option 2: Empty grid that can be populated later
            grid = pres.add_autogrid(slide, None, rows=2, cols=3)
            grid.add_to_cell(0, 0, slide.add_text, text="Cell 0,0")
            grid.add_to_cell(1, 2, slide.add_text, text="Cell 1,2")
            ```
        """
        from easypptx.grid import Grid

        # If content_funcs is None and rows/cols are provided, create an empty grid
        if content_funcs is None:
            # Make sure rows and cols are specified for empty grid
            if rows is None or cols is None:
                rows = rows or 1
                cols = cols or 1

            # Adjust grid position and dimensions if a title is provided
            adjusted_y = y
            adjusted_height = height

            if title:
                # Reserve the title band above the grid
                adjusted_y, adjusted_height = shift_band(y, height, title_height, self._slide_height_emu)

                # Add the title to the slide
                slide.add_text(
                    text=title,
                    x=x,
                    y=y,
                    width=width,
                    height=title_height,
                    font_size=24,
                    font_bold=True,
                    align=title_align,
                    vertical="middle",
                )

            # Create empty grid with specified dimensions
            grid = Grid(
                parent=slide,
                x=x,
                y=adjusted_y,
                width=width,
                height=adjusted_height,
                rows=rows,
                cols=cols,
                padding=padding,
            )
        else:
            # Use the Grid.autogrid method for content_funcs
            grid = Grid.autogrid(
                parent=slide,
                content_funcs=content_funcs,
                rows=rows,
                cols=cols,
                x=x,
                y=y,
                width=width,
                height=height,
                padding=padding,
                title=title,
                title_height=title_height,
                column_major=column_major,
            )

        return grid

    def add_grid_slide(
        self,
        rows: int,
        cols: int,
        title: str | None = None,
        subtitle: str | None = None,
        title_height: float | str | None = None,  # None means use template default
        subtitle_height: float | str | None = None,  # None means use template default
        x: float | str | None = None,  # None means use template default
        y: float | str | None = None,  # None means use template default
        width: float | str | None = None,  # None means use template default
        height: float | str | None = None,  # None means use template default
        padding: float | None = None,  # None means use template default
        bg_color: str | tuple[int, int, int] | None = None,  # None means use template default
        title_font_size: int | None = None,  # None means use template default
        subtitle_font_size: int | None = None,  # None means use template default
        title_align: str | None = None,  # None means use template default
        subtitle_align: str | None = None,  # None means use template default
        title_padding: str | float | None = None,  # None means use template default
        title_x_padding: str | float | None = None,  # None means use template default
        title_y_padding: str | float | None = None,  # None means use template default
        subtitle_padding: str | float | None = None,  # None means use template default
        subtitle_x_padding: str | float | None = None,  # None means use template default
        subtitle_y_padding: str | float | None = None,  # None means use template default
        content_padding: str | float | None = None,  # None means use template default
        content_x_padding: str | float | None = None,  # None means use template default
        content_y_padding: str | float | None = None,  # None means use template default
    ) -> tuple[Slide, Grid]:
        """Add a slide with a grid layout.

        This method creates a new slide with an empty grid that can be populated later.
        It provides flexible options for positioning and sizing the grid, as well as
        adding a title and subtitle to the slide.

        All parameters can be set to None to use the template defaults if a template is
        registered with the presentation.

        Args:
            rows: Number of rows in the grid
            cols: Number of columns in the grid
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            title_height: Height of the title area (default: template default or "10%")
            subtitle_height: Height of the subtitle area (default: template default or "5%")
            x: X position of the grid as percentage or absolute value (default: template default or "0%")
            y: Y position of the grid as percentage or absolute value (default: template default or "0%")
            width: Width of the grid as percentage or absolute value (default: template default or "100%")
            height: Height of the grid as percentage or absolute value (default: template default or "100%")
            padding: Padding between cells as percentage of cell size (default: template default or 5.0)
            bg_color: Background color for the slide (default: template default or None)
            title_font_size: Font size for the title (default: template default or 24)
            subtitle_font_size: Font size for the subtitle (default: template default or 18)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: template default or "center")
            subtitle_align: Text alignment for the subtitle, one of "left", "center", "right" (default: template default or "center")
            title_padding: Padding around the title, applies to both x and y (default: template default or None)
            title_x_padding: Horizontal padding for title, overridden by title_padding if provided (default: template default or None)
            title_y_padding: Vertical padding for title, overridden by title_padding if provided (default: template default or None)
            subtitle_padding: Padding around the subtitle, applies to both x and y (default: template default or None)
            subtitle_x_padding: Horizontal padding for subtitle, overridden by subtitle_padding if provided (default: template default or None)
            subtitle_y_padding: Vertical padding for subtitle, overridden by subtitle_padding if provided (default: template default or None)
            content_padding: Padding around the content area, applies to both x and y (default: template default or None)
            content_x_padding: Horizontal padding for content, overridden by content_padding if provided (default: template default or None)
            content_y_padding: Vertical padding for content, overridden by content_padding if provided (default: template default or None)

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            # Create a slide with a 3x2 grid and a title
            slide, grid = pres.add_grid_slide(
                rows=3,
                cols=2,
                title="Features Overview",
                subtitle="Product Capabilities",
                padding=5.0,
                title_align="left",
                title_padding="5%"
            )

            # Add content to specific cells
            grid[0, 0].add_text("Feature 1", font_bold=True)
            grid[0, 1].add_image("image1.png")

            # Access cells using flat indexing (row-major order)
            grid[2].add_text("Feature 2", font_bold=True)  # Third cell (flat index 2)
            grid[3].add_text("Description of Feature 2")   # Fourth cell (flat index 3)
            ```
        """
        # Get template defaults if available
        template_defaults = {}
        if self._default_template is not None:
            try:
                template = self.template_manager.get(self._default_template)
                if "defaults" in template and "grid_slide" in template["defaults"]:
                    template_defaults = template["defaults"]["grid_slide"]
            except (ValueError, KeyError):
                # If template lookup fails, use empty defaults
                pass

        # Apply defaults for parameters that are None
        x_val = x if x is not None else template_defaults.get("x", "0%")
        y_val = y if y is not None else template_defaults.get("y", "0%")
        width_val = width if width is not None else template_defaults.get("width", "100%")
        height_val = height if height is not None else template_defaults.get("height", "100%")
        padding_val = padding if padding is not None else template_defaults.get("padding", 5.0)
        title_height_val = title_height if title_height is not None else template_defaults.get("title_height", "10%")
        subtitle_height_val = (
            subtitle_height if subtitle_height is not None else template_defaults.get("subtitle_height", "5%")
        )
        title_font_size_val = (
            title_font_size if title_font_size is not None else template_defaults.get("title_font_size", 24)
        )
        subtitle_font_size_val = (
            subtitle_font_size if subtitle_font_size is not None else template_defaults.get("subtitle_font_size", 18)
        )
        bg_color_val = bg_color if bg_color is not None else template_defaults.get("bg_color")

        # Create a new slide
        slide = self.add_slide(bg_color=bg_color_val)
        slide_height_emu = self._slide_height_emu

        # Calculate positions and dimensions
        adjusted_y = y_val
        adjusted_height = height_val

        # Apply content padding if specified (for grid positioning)
        grid_x = x_val
        content_padding_val = (
            content_padding if content_padding is not None else template_defaults.get("content_padding")
        )
        content_x_padding_val = (
            content_x_padding if content_x_padding is not None else template_defaults.get("content_x_padding")
        )

        if content_padding_val is not None or content_x_padding_val is not None:
            content_x = content_padding_val if content_padding_val is not None else content_x_padding_val
            if content_x is not None:
                grid_x = content_x

        # Add title if provided
        if title:
            # Calculate title position with padding
            title_x = x_val
            title_y = y_val

            # Get title padding values from template defaults if not specified
            title_padding_val = title_padding if title_padding is not None else template_defaults.get("title_padding")
            title_x_padding_val = (
                title_x_padding if title_x_padding is not None else template_defaults.get("title_x_padding")
            )
            title_y_padding_val = (
                title_y_padding if title_y_padding is not None else template_defaults.get("title_y_padding")
            )

            # Apply title padding if specified
            if title_padding_val is not None or title_x_padding_val is not None:
                title_x_padding_value = title_padding_val if title_padding_val is not None else title_x_padding_val
                if title_x_padding_value is not None:
                    title_x = title_x_padding_value

            if title_padding_val is not None or title_y_padding_val is not None:
                title_y_padding_value = title_padding_val if title_padding_val is not None else title_y_padding_val
                if title_y_padding_value is not None:
                    title_y = title_y_padding_value

            # Get title alignment
            title_align_val = title_align
            if title_align_val is None:
                # First check template defaults for grid_slide
                title_align_val = template_defaults.get("title_align")

                # If still None and we have a default template, try to get more defaults
                if title_align_val is None and self._default_template is not None:
                    try:
                        # Try to get template settings
                        template_data = self.template_manager.get(self._default_template)

                        # Check global defaults if available
                        if "defaults" in template_data and "global" in template_data["defaults"]:
                            # Check for title_align in global defaults
                            title_align_val = template_data["defaults"]["global"].get("title_align")

                            # If still None, check for generic align in global defaults
                            if title_align_val is None:
                                title_align_val = template_data["defaults"]["global"].get("align")

                        # If still None, check dedicated title section alignment
                        if title_align_val is None and "title" in template_data and "align" in template_data["title"]:
                            title_align_val = template_data["title"]["align"]
                    except (ValueError, KeyError):
                        pass

            # Default to center if still None
            if title_align_val is None:
                title_align_val = "center"

            # Get title font properties from template if available
            title_font_bold = template_defaults.get("title_font_bold", True)
            title_vertical = template_defaults.get("title_vertical", "middle")

            # Add the title to the slide
            slide.add_text(
                text=title,
                x=title_x,
                y=title_y,
                width=width_val,
                height=title_height_val,
                font_size=title_font_size_val,
                font_bold=title_font_bold,
                align=title_align_val,
                vertical=title_vertical,
            )

            # Reserve the title band above the grid
            adjusted_y, adjusted_height = shift_band(y_val, height_val, title_height_val, slide_height_emu)

        # Add subtitle if provided
        if subtitle:
            # Calculate subtitle position with padding
            subtitle_x = x_val
            subtitle_y = adjusted_y

            # Get subtitle padding values from template defaults if not specified
            subtitle_padding_val = (
                subtitle_padding if subtitle_padding is not None else template_defaults.get("subtitle_padding")
            )
            subtitle_x_padding_val = (
                subtitle_x_padding if subtitle_x_padding is not None else template_defaults.get("subtitle_x_padding")
            )
            subtitle_y_padding_val = (
                subtitle_y_padding if subtitle_y_padding is not None else template_defaults.get("subtitle_y_padding")
            )

            # Apply subtitle padding if specified
            if subtitle_padding_val is not None or subtitle_x_padding_val is not None:
                subtitle_x_padding_value = (
                    subtitle_padding_val if subtitle_padding_val is not None else subtitle_x_padding_val
                )
                if subtitle_x_padding_value is not None:
                    subtitle_x = subtitle_x_padding_value

            if subtitle_padding_val is not None or subtitle_y_padding_val is not None:
                subtitle_y_padding_value = (
                    subtitle_padding_val if subtitle_padding_val is not None else subtitle_y_padding_val
                )
                subtitle_y = subtitle_y_padding_value if subtitle_y_padding_value is not None else adjusted_y

            # Get subtitle alignment
            subtitle_align_val = subtitle_align
            if subtitle_align_val is None:
                # First check template defaults for grid_slide
                subtitle_align_val = template_defaults.get("subtitle_align")

                # If still None and we have a default template, try to get more defaults
                if subtitle_align_val is None and self._default_template is not None:
                    try:
                        # Try to get template settings
                        template_data = self.template_manager.get(self._default_template)

                        # Check global defaults if available
                        if "defaults" in template_data and "global" in template_data["defaults"]:
                            # Check for subtitle_align in global defaults
                            subtitle_align_val = template_data["defaults"]["global"].get("subtitle_align")

                            # If still None, check for generic align in global defaults
                            if subtitle_align_val is None:
                                subtitle_align_val = template_data["defaults"]["global"].get("align")

                        # If still None, check dedicated subtitle section alignment
                        if (
                            subtitle_align_val is None
                            and "subtitle" in template_data
                            and "align" in template_data["subtitle"]
                        ):
                            subtitle_align_val = template_data["subtitle"]["align"]
                    except (ValueError, KeyError):
                        pass

            # Default to center if still None
            if subtitle_align_val is None:
                subtitle_align_val = "center"

            # Get subtitle font properties from template if available
            subtitle_font_bold = template_defaults.get("subtitle_font_bold", False)
            subtitle_vertical = template_defaults.get("subtitle_vertical", "middle")

            # Add the subtitle to the slide
            slide.add_text(
                text=subtitle,
                x=subtitle_x,
                y=subtitle_y,
                width=width_val,
                height=subtitle_height_val,
                font_size=subtitle_font_size_val,
                font_bold=subtitle_font_bold,
                align=subtitle_align_val,
                vertical=subtitle_vertical,
            )

            # Reserve the subtitle band above the grid
            adjusted_y, adjusted_height = shift_band(adjusted_y, adjusted_height, subtitle_height_val, slide_height_emu)

        # Apply content y padding if specified (explicit parameters win over
        # template defaults; the padding is applied exactly once)
        content_y_padding_val = (
            content_y_padding if content_y_padding is not None else template_defaults.get("content_y_padding")
        )
        grid_y, adjusted_height = apply_content_padding(
            adjusted_y, adjusted_height, content_padding_val, content_y_padding_val, slide_height_emu
        )

        # Create the grid
        grid = Grid(
            parent=slide,
            x=grid_x,
            y=grid_y,
            width=width_val,
            height=adjusted_height,
            rows=rows,
            cols=cols,
            padding=padding_val,
        )

        # Apply template defaults from template to grid if available
        if self._default_template is not None:
            try:
                template = self.template_manager.get(self._default_template)
                grid.apply_template_defaults(template)
            except (ValueError, KeyError):
                # If template lookup fails, proceed without template defaults
                pass

        return slide, grid

    def add_autogrid_slide(
        self,
        content_funcs: list | None = None,
        rows: int | None = None,
        cols: int | None = None,
        title: str | None = None,
        title_height: float | str = "10%",
        padding: float = 5.0,
        bg_color: str | tuple[int, int, int] | None = None,
        title_align: str = "center",
        column_major: bool = True,  # Use column-major order by default
        title_padding: str | float | None = None,
        title_x_padding: str | float | None = None,
        title_y_padding: str | float | None = None,
        title_font_size: int = 24,
        content_padding: str | float | None = None,
        content_x_padding: str | float | None = None,
        content_y_padding: str | float | None = None,
    ) -> tuple[Slide, Grid]:
        """Add a slide with an autogrid layout.

        This method creates a new slide and automatically places the provided
        content functions into a grid. If content_funcs is None, it creates an
        empty grid with the specified rows and columns that can be populated later.

        Args:
            content_funcs: List of functions that add content to the slide, or None for an empty grid
            rows: Number of rows (default: None, calculated automatically when content_funcs provided)
            cols: Number of columns (default: None, calculated automatically when content_funcs provided)
            title: Optional title for the slide (default: None)
            title_height: Height of the title area (default: "10%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            bg_color: Background color for the slide (default: None)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            column_major: Whether to fill cells in column-major order (default: True)
                         When True, fills cells down columns first, resulting in a visual layout
                         that matches the specified rows and columns when content is added sequentially.
                         When False, fills cells across rows first.
            title_padding: Padding around the title, applies to both x and y (default: None)
            title_x_padding: Horizontal padding for title, overridden by title_padding if provided (default: None)
            title_y_padding: Vertical padding for title, overridden by title_padding if provided (default: None)
            title_font_size: Font size for the title (default: 24)
            content_padding: Padding around the content area, applies to both x and y (default: None)
            content_x_padding: Horizontal padding for content, overridden by content_padding if provided (default: None)
            content_y_padding: Vertical padding for content, overridden by content_padding if provided (default: None)

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            # With content functions
            def create_text1():
                return slide.add_text("Text 1")

            def create_text2():
                return slide.add_text("Text 2")

            content_funcs = [create_text1, create_text2]
            slide, grid = pres.add_autogrid_slide(
                content_funcs,
                title="Auto Grid Slide",
                title_align="left",
                title_padding="5%",
                content_padding="2%"
            )

            # With empty grid
            slide, grid = pres.add_autogrid_slide(None, rows=4, cols=2, title="Features")

            # Add content directly to rows
            grid[0].add_text("Feature 1", font_bold=True)
            grid[0].add_text("Description 1")
            grid[1].add_text("Feature 2", font_bold=True)
            grid[1].add_text("Description 2")
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)
        slide_height_emu = self._slide_height_emu

        # If content_funcs is None and rows/cols are provided, ensure they have values
        if content_funcs is None and (rows is None or cols is None):
            rows = rows or 1
            cols = cols or 1

        # Set default grid position and dimensions
        grid_y: str | float = "0%"
        grid_height: str | float = "100%"

        # Apply content x padding if specified (for grid positioning)
        content_x = resolve_padding(content_padding, content_x_padding)
        grid_x = content_x if content_x is not None else "0%"

        if title:
            # Calculate title position with padding
            title_x = resolve_padding(title_padding, title_x_padding)
            title_y = resolve_padding(title_padding, title_y_padding)

            # Add the title to the slide
            slide.add_text(
                text=title,
                x=title_x if title_x is not None else "0%",
                y=title_y if title_y is not None else "0%",
                width="100%",
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Reserve the title band above the grid
            grid_y, grid_height = shift_band("0%", "100%", title_height, slide_height_emu)

        # Apply content y padding if specified
        grid_y, grid_height = apply_content_padding(
            grid_y, grid_height, content_padding, content_y_padding, slide_height_emu
        )

        # Create the autogrid without its own title (added above if requested)
        grid = self.add_autogrid(
            slide=slide,
            content_funcs=content_funcs,
            rows=rows,
            cols=cols,
            x=grid_x,
            y=grid_y,
            width="100%",
            height=grid_height,
            padding=padding,
            title=None,
            title_align=title_align,
            column_major=column_major,
        )

        # Apply template defaults to grid if available
        if self._default_template is not None:
            try:
                template = self.template_manager.get(self._default_template)
                grid.apply_template_defaults(template)
            except (ValueError, KeyError):
                # If template lookup fails, proceed without template defaults
                pass

        return slide, grid

    def add_pyplot_slide(
        self,
        figure,
        title: str | None = None,
        subtitle: str | None = None,
        label: str | None = None,
        x: float | str = "10%",
        y: float | str = "20%",
        width: float | str = "80%",
        height: float | str = "70%",
        title_height: float | str = "10%",
        subtitle_height: float | str = "5%",
        dpi: int = 300,
        file_format: str = "png",
        bg_color: str | tuple[int, int, int] | None = None,
        title_font_size: int = 24,
        subtitle_font_size: int = 18,
        label_font_size: int = 14,
        border: bool = False,
        border_color: str = "black",
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
        title_align: str = "center",
        subtitle_align: str = "center",
        label_align: str = "center",
        title_padding: str | float | None = None,
        title_x_padding: str | float | None = None,
        title_y_padding: str | float | None = None,
        subtitle_padding: str | float | None = None,
        subtitle_x_padding: str | float | None = None,
        subtitle_y_padding: str | float | None = None,
        content_padding: str | float | None = None,
        content_x_padding: str | float | None = None,
        content_y_padding: str | float | None = None,
        label_padding: str | float | None = None,
        label_x_padding: str | float | None = None,
        label_y_padding: str | float | None = "1%",
    ) -> tuple[Slide, PPTXShape]:
        """Add a slide with a title and a matplotlib/seaborn figure.

        This method provides more flexibility than add_matplotlib_slide and add_seaborn_slide
        by allowing control over positioning, titles, and styling, similar to add_grid_slide.

        Args:
            figure: Matplotlib or Seaborn figure object
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            label: Optional caption for the figure, displayed below (default: None)
            x: X position of the figure as percentage or absolute value (default: "10%")
            y: Y position of the figure as percentage or absolute value (default: "20%")
            width: Width of the figure as percentage or absolute value (default: "80%")
            height: Height of the figure as percentage or absolute value (default: "70%")
            title_height: Height of the title area (default: "10%")
            subtitle_height: Height of the subtitle area (default: "5%")
            dpi: Resolution for the figure in dots per inch (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            bg_color: Background color for the slide (default: None)
            title_font_size: Font size for the title (default: 24)
            subtitle_font_size: Font size for the subtitle (default: 18)
            label_font_size: Font size for the caption (default: 14)
            border: Whether to add a border around the figure (default: False)
            border_color: Color for the border (default: "black")
            shadow: Whether to add a shadow effect to the figure (default: False)
            maintain_aspect_ratio: Whether to maintain the figure's aspect ratio (default: True)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            subtitle_align: Text alignment for the subtitle, one of "left", "center", "right" (default: "center")
            label_align: Text alignment for the caption, one of "left", "center", "right" (default: "center")
            title_padding: Padding around the title, applies to both x and y (default: None)
            title_x_padding: Horizontal padding for title, overridden by title_padding if provided (default: None)
            title_y_padding: Vertical padding for title, overridden by title_padding if provided (default: None)
            subtitle_padding: Padding around the subtitle, applies to both x and y (default: None)
            subtitle_x_padding: Horizontal padding for subtitle, overridden by subtitle_padding if provided (default: None)
            subtitle_y_padding: Vertical padding for subtitle, overridden by subtitle_padding if provided (default: None)
            content_padding: Padding around the content area, applies to both x and y (default: None)
            content_x_padding: Horizontal padding for content, overridden by content_padding if provided (default: None)
            content_y_padding: Vertical padding for content, overridden by content_padding if provided (default: None)
            label_padding: Padding around the label, applies to both x and y (default: None)
            label_x_padding: Horizontal padding for label, overridden by label_padding if provided (default: None)
            label_y_padding: Vertical padding between content and label (default: "1%")

        Returns:
            A tuple containing (Slide, PPTXShape) where PPTXShape is the figure shape

        Example:
            ```python
            import matplotlib.pyplot as plt

            # Create a matplotlib figure
            plt.figure(figsize=(8, 6))
            plt.plot([1, 2, 3, 4], [1, 4, 9, 16])
            plt.title('Sample Plot')
            plt.grid(True)

            # Add it to a presentation with title and subtitle
            slide, pyplot = pres.add_pyplot_slide(
                figure=plt.gcf(),
                title="Data Visualization",
                subtitle="Matplotlib Example",
                label="Figure 1: Sample Plot",
                dpi=300,
                title_align="left",
                title_padding="5%",
                content_padding="2%"
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)
        slide_height_emu = self._slide_height_emu

        # Calculate positions and dimensions
        adjusted_y = y
        adjusted_height = height

        # Determine content x-padding
        content_x_pad = resolve_padding(content_padding, content_x_padding)
        content_x_val = content_x_pad if content_x_pad is not None else x

        # Add title if provided
        if title:
            title_x = resolve_padding(title_padding, title_x_padding)
            title_y = resolve_padding(title_padding, title_y_padding)

            slide.add_text(
                text=title,
                x=title_x if title_x is not None else "0%",
                y=title_y if title_y is not None else "0%",
                width="100%",
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Reserve the title band above the figure
            adjusted_y, adjusted_height = shift_band(y, height, title_height, slide_height_emu)

        # Add subtitle if provided
        if subtitle:
            subtitle_x = resolve_padding(subtitle_padding, subtitle_x_padding)
            subtitle_y = resolve_padding(subtitle_padding, subtitle_y_padding)

            slide.add_text(
                text=subtitle,
                x=subtitle_x if subtitle_x is not None else "0%",
                y=subtitle_y if subtitle_y is not None else adjusted_y,
                width="100%",
                height=subtitle_height,
                font_size=subtitle_font_size,
                align=subtitle_align,
                vertical="middle",
            )

            # Reserve the subtitle band above the figure
            adjusted_y, adjusted_height = shift_band(adjusted_y, adjusted_height, subtitle_height, slide_height_emu)

        # Apply content y padding if specified
        figure_y, adjusted_height = apply_content_padding(
            adjusted_y, adjusted_height, content_padding, content_y_padding, slide_height_emu
        )

        # Create a style dictionary for the pyplot
        style = {
            "border": border,
            "border_color": border_color,
            "shadow": shadow,
            "maintain_aspect_ratio": maintain_aspect_ratio,
        }

        # Add the pyplot to the slide
        pyplot = Pyplot.add(
            slide=slide,
            figure=figure,
            position={
                "x": content_x_val,
                "y": figure_y,
                "width": width,
                "height": adjusted_height,
            },
            dpi=dpi,
            file_format=file_format,
            style=style,
        )

        # Add label if specified
        if label:
            label_x = resolve_padding(label_padding, label_x_padding)
            label_y_pad = resolve_padding(label_padding, label_y_padding)

            # Place the label just below the figure
            label_y = pct(
                to_percent(figure_y, slide_height_emu)
                + to_percent(adjusted_height, slide_height_emu)
                + (to_percent(label_y_pad, slide_height_emu) if label_y_pad is not None else 1.0)
            )

            slide.add_text(
                text=label,
                x=label_x if label_x is not None else "0%",
                y=label_y,
                width="100%",
                height="5%",
                font_size=label_font_size,
                align=label_align,
                vertical="top",
            )

        return slide, pyplot
