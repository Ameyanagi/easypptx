"""Slide module for EasyPPTX."""

from __future__ import annotations

import math
from typing import IO, TYPE_CHECKING, Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.autoshape import Shape as PPTXShape
from pptx.slide import Slide as PPTXSlide
from pptx.util import Inches, Pt

from easypptx.common import (
    ALIGN,
    COLORS,
    DEFAULT_FONT,
    VERTICAL,
    is_dataframe,
    merge_defaults,
    normalize_color,
    resolve_color,
    warn_ignored_kwargs,
)
from easypptx.positioning import PositionType, to_inches
from easypptx.styles import ChartStyle, TableStyle, TextStyle

if TYPE_CHECKING:
    from pptx.chart.chart import Chart as PPTXChart
    from pptx.table import Table as PPTXTable

# Default slide dimensions: standard 16:9 (13.33 x 7.5 inches) in EMUs
_DEFAULT_SLIDE_WIDTH = 12192768
_DEFAULT_SLIDE_HEIGHT = 6858000


class Slide:
    """Class representing a slide in a PowerPoint presentation.

    This class provides methods for adding and manipulating content on a slide,
    such as text, images, tables, charts, and shapes.

    Attributes:
        pptx_slide: The underlying python-pptx Slide object

    Examples:
        ```python
        # Add text to a slide using inches
        slide.add_text("Hello World", x=2, y=2)

        # Add text to a slide using percentages
        slide.add_text("Hello World", x="20%", y="30%")

        # Add an image using percentages
        slide.add_image("image.png", x="10%", y="20%", width="50%", height="30%")

        # Add a table
        slide.add_table([["Name", "Value"], ["A", 1]], x="10%", y="20%")
        ```
    """

    def __init__(
        self,
        pptx_slide: PPTXSlide,
        slide_width: int | None = None,
        slide_height: int | None = None,
    ) -> None:
        """Initialize a Slide object.

        Args:
            pptx_slide: The python-pptx Slide object
            slide_width: Slide width in EMUs (default: None, discovered from the slide)
            slide_height: Slide height in EMUs (default: None, discovered from the slide)
        """
        self.pptx_slide = pptx_slide
        self.user_data: dict[str, Any] = {}

        # Store template defaults when applied from a template
        self.template_defaults: dict[str, dict[str, Any]] = {
            "text": {},
            "image": {},
            "shape": {},
            "table": {},
            "chart": {},
            "global": {},
        }

        # Cache slide dimensions to avoid recalculating them
        self._slide_width = slide_width if slide_width is not None else self._discover_dimension("slide_width")
        self._slide_height = slide_height if slide_height is not None else self._discover_dimension("slide_height")

    def apply_template_defaults(self, template_data: dict[str, Any]) -> None:
        """Apply template defaults to this slide.

        This method extracts default method arguments from a template and stores them
        for later use by the add_xxx methods.

        Args:
            template_data: Template data dictionary
        """
        defaults = template_data.get("defaults", {})
        for element_type in self.template_defaults:
            if element_type in defaults:
                self.template_defaults[element_type] = defaults[element_type]

    def merge_with_defaults(self, method_type: str, kwargs: dict[str, Any]) -> dict[str, Any]:
        """Merge provided arguments with template defaults.

        Args:
            method_type: The type of method ("text", "image", etc.)
            kwargs: Keyword arguments provided to the method

        Returns:
            Dictionary with merged arguments, where provided args override defaults
        """
        return merge_defaults(self.template_defaults, method_type, kwargs)

    def _convert_position(self, value: PositionType, slide_dimension: int, is_width: bool = True) -> float:
        """Convert a position value (percentage string or inches) to inches.

        Args:
            value: Position value (percentage string like "20%" or absolute inches)
            slide_dimension: The total slide dimension (width or height) in EMUs
            is_width: Unused, kept for backward compatibility

        Returns:
            Position value in inches
        """
        return to_inches(value, slide_dimension)

    def _discover_dimension(self, attribute: str) -> int:
        """Discover a slide dimension (width or height) in EMUs.

        Falls back to standard 16:9 dimensions when the presentation is not
        reachable from the slide (e.g. detached slides or test doubles).

        Args:
            attribute: "slide_width" or "slide_height"

        Returns:
            The dimension in English Metric Units (EMUs)
        """
        default = _DEFAULT_SLIDE_WIDTH if attribute == "slide_width" else _DEFAULT_SLIDE_HEIGHT
        try:
            value = getattr(self.pptx_slide.part.package.presentation_part.presentation, attribute)
        except (AttributeError, TypeError):
            return default
        return value if isinstance(value, int) else default

    def _get_slide_width(self) -> int:
        """Get the slide width in EMUs."""
        return self._slide_width

    def _get_slide_height(self) -> int:
        """Get the slide height in EMUs."""
        return self._slide_height

    def _template_value(self, method_type: str, key: str, provided: Any, fallback: Any) -> Any:
        """Resolve a parameter: explicit value, else template default, else fallback."""
        if provided is not None:
            return provided
        method_defaults = self.template_defaults.get(method_type, {})
        global_defaults = self.template_defaults.get("global", {})
        return method_defaults.get(key, global_defaults.get(key, fallback))

    def add_text(
        self,
        text: str,
        x: PositionType | None = None,
        y: PositionType | None = None,
        width: PositionType | None = None,
        height: PositionType | None = None,
        font_size: int | None = None,
        font_bold: bool | None = None,
        font_italic: bool | None = None,
        font_name: str | None = None,
        align: str | None = None,
        vertical: str | None = None,
        color: str | tuple[int, int, int] | None = None,
        style: TextStyle | None = None,
        **kwargs: Any,
    ) -> PPTXShape:
        """Add a text box to the slide.

        Args:
            text: The text content
            x: X position as percent (number or "N%") or in_() length (default: 5)
            y: Y position as percent or in_() length (default: 5)
            width: Width as percent or in_() length (default: 90)
            height: Height as percent or in_() length (default: 10)
            font_size: Font size in points (default: 18)
            font_bold: Whether text should be bold (default: False)
            font_italic: Whether text should be italic (default: False)
            font_name: Font name (default: "Meiryo")
            align: Text alignment, one of "left", "center", "right" (default: "left")
            vertical: Vertical alignment, one of "top", "middle", "bottom" (default: "top")
            color: Text color as string name from COLORS dict or RGB tuple (default: None)
            **kwargs: "vertical_align" is accepted as an alias for "vertical";
                any other unknown parameter triggers a warning

        Returns:
            The created shape object
        """
        # Accept the alias used elsewhere in the library
        if vertical is None and "vertical_align" in kwargs:
            vertical = kwargs.pop("vertical_align")
        warn_ignored_kwargs("Slide.add_text", kwargs)

        # A TextStyle fills any formatting the caller left unset
        if style is not None:
            sk = style.to_kwargs()
            font_size = font_size if font_size is not None else sk.get("font_size")
            font_bold = font_bold if font_bold is not None else sk.get("font_bold")
            font_italic = font_italic if font_italic is not None else sk.get("font_italic")
            font_name = font_name if font_name is not None else sk.get("font_name")
            align = align if align is not None else sk.get("align")
            vertical = vertical if vertical is not None else sk.get("vertical")
            color = color if color is not None else sk.get("color")

        # Resolve parameters: explicit value > template default > fallback
        x_val = self._template_value("text", "x", x, 5)
        y_val = self._template_value("text", "y", y, 5)
        width_val = self._template_value("text", "width", width, 90)
        height_val = self._template_value("text", "height", height, 10)
        font_size_val = self._template_value("text", "font_size", font_size, 18)
        font_bold_val = self._template_value("text", "font_bold", font_bold, False)
        font_italic_val = self._template_value("text", "font_italic", font_italic, False)
        font_name_val = self._template_value("text", "font_name", font_name, DEFAULT_FONT)
        align_val = self._template_value("text", "align", align, "left")
        vertical_val = self._template_value("text", "vertical", vertical, "top")
        color_val = normalize_color(self._template_value("text", "color", color, None))

        # Convert position values to inches
        x_inches = self._convert_position(x_val, self._slide_width)
        y_inches = self._convert_position(y_val, self._slide_height)
        width_inches = self._convert_position(width_val, self._slide_width)
        height_inches = self._convert_position(height_val, self._slide_height)

        # Create the textbox
        text_box = self.pptx_slide.shapes.add_textbox(
            Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True  # Enable word wrap for better text display
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Set vertical alignment
        if vertical_val in VERTICAL:
            text_frame.vertical_anchor = VERTICAL[vertical_val]

        # Apply text formatting
        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size_val)
        p.font.bold = font_bold_val
        p.font.italic = font_italic_val
        p.font.name = font_name_val

        # Set horizontal alignment
        if align_val in ALIGN:
            p.alignment = ALIGN[align_val]

        # Set text color
        rgb = resolve_color(color_val)
        if rgb is not None:
            p.font.color.rgb = rgb

        return text_box

    def add_image(
        self,
        image_path: str | IO[bytes],
        x: PositionType = 5,
        y: PositionType = 5,
        width: PositionType | None = None,
        height: PositionType | None = None,
        **kwargs: Any,
    ) -> PPTXShape:
        """Add an image to the slide.

        Args:
            image_path: Path to the image file, or a binary file-like object
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: None, maintains aspect ratio)
            height: Height in inches or percentage (default: None, maintains aspect ratio)
            **kwargs: Unknown parameters trigger a warning

        Returns:
            The created picture shape

        Raises:
            FileNotFoundError: If the image file doesn't exist
        """
        # maintain_aspect_ratio is implemented by Image.add, which sizes the
        # missing dimension before calling this method; here it is a no-op.
        kwargs.pop("maintain_aspect_ratio", None)
        warn_ignored_kwargs("Slide.add_image", kwargs)

        x_inches = self._convert_position(x, self._slide_width)
        y_inches = self._convert_position(y, self._slide_height)

        width_inches = None
        height_inches = None
        if width is not None:
            width_inches = Inches(self._convert_position(width, self._slide_width))
        if height is not None:
            height_inches = Inches(self._convert_position(height, self._slide_height))

        picture: Any = self.pptx_slide.shapes.add_picture(
            image_path, Inches(x_inches), Inches(y_inches), width_inches, height_inches
        )
        return picture

    @property
    def shapes(self) -> list[Any]:
        """Get all shapes on the slide.

        Returns:
            List of shape objects
        """
        return list(self.pptx_slide.shapes)

    def clear(self) -> None:
        """Remove all shapes from the slide."""
        for shape in self.pptx_slide.shapes:
            self.pptx_slide.shapes._spTree.remove(shape._element)

    def add_shape(
        self,
        shape_type: MSO_SHAPE | str = MSO_SHAPE.RECTANGLE,
        x: PositionType = 5,
        y: PositionType = 5,
        width: PositionType = 40,
        height: PositionType = 10,
        fill_color: str | tuple[int, int, int] | None = None,
        **kwargs: Any,
    ) -> PPTXShape:
        """Add a shape to the slide.

        Args:
            shape_type: The shape type, as an MSO_SHAPE value or its name,
                e.g. "ROUNDED_RECTANGLE" (default: MSO_SHAPE.RECTANGLE)
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 5.0)
            height: Height in inches or percentage (default: 1.0)
            fill_color: Fill color as string name from COLORS dict or RGB tuple (default: None)
            **kwargs: Unknown parameters trigger a warning

        Returns:
            The created shape object

        Raises:
            ValueError: If shape_type is a string that is not an MSO_SHAPE name
        """
        warn_ignored_kwargs("Slide.add_shape", kwargs)

        if isinstance(shape_type, str):
            name = shape_type.upper()
            try:
                shape_type = getattr(MSO_SHAPE, name)
            except AttributeError:
                raise ValueError(f"Unknown shape type: {shape_type!r} (use an MSO_SHAPE name)") from None

        x_inches = self._convert_position(x, self._slide_width)
        y_inches = self._convert_position(y, self._slide_height)
        width_inches = self._convert_position(width, self._slide_width)
        height_inches = self._convert_position(height, self._slide_height)

        shape = self.pptx_slide.shapes.add_shape(
            shape_type, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )

        rgb = resolve_color(fill_color)
        if rgb is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb

        return shape

    def add_table(
        self,
        data: Any,
        x: PositionType = 5,
        y: PositionType = 20,
        width: PositionType | None = None,
        height: PositionType | None = None,
        has_header: bool | None = None,
        style: int | dict | TableStyle | None = None,
        **kwargs: Any,
    ) -> PPTXTable:
        """Add a table to the slide.

        Args:
            data: Table data as a list of lists or pandas DataFrame
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Total width in inches or percentage (default: None, auto-sized)
            height: Total height in inches or percentage (default: None, auto-sized)
            has_header: Whether the first row is a header (default: True)
            style: Table style ID, or a style dict with a "style_id" key (default: None)
            **kwargs: "first_row_header" is accepted as an alias for "has_header";
                any other unknown parameter triggers a warning

        Returns:
            The created table object
        """
        from easypptx.table import Table

        if "first_row_header" in kwargs:
            has_header = kwargs.pop("first_row_header")
        warn_ignored_kwargs("Slide.add_table", kwargs)

        if isinstance(style, TableStyle):
            if has_header is None:
                has_header = style.has_header
            style_id = style.style_id
        elif isinstance(style, dict):
            style_id = style.get("style_id")
        else:
            style_id = style
        if has_header is None:
            has_header = True

        table_data = [list(data.columns), *data.values.tolist()] if is_dataframe(data) else data

        return Table(self).add(
            data=table_data,
            x=x,
            y=y,
            width=width,
            height=height,
            first_row_header=has_header,
            style=style_id,
        )

    def add_chart(
        self,
        data: Any = None,
        chart_type: str | None = None,
        x: PositionType = 10,
        y: PositionType = 20,
        width: PositionType = 60,
        height: PositionType = 60,
        categories: list | None = None,
        values: list | None = None,
        category_column: str | int | None = None,
        value_columns: str | list[str] | int | list[int] | None = None,
        title: str | None = None,
        has_legend: bool | None = None,
        legend_position: str | None = None,
        style: ChartStyle | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a chart to the slide.

        Either pass explicit ``categories`` and ``values`` lists, or pass
        ``data`` (a list of lists with a header row, or a pandas DataFrame)
        together with optional ``category_column`` / ``value_columns``.

        Args:
            data: Chart data as a list of lists or pandas DataFrame (default: None)
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', 'area', 'scatter')
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            categories: Explicit list of category labels (default: None)
            values: Explicit list of data values (default: None)
            category_column: Name or index of the column to use as categories (default: None)
            value_columns: Name(s) or index(es) of column(s) to use as values (default: None)
            title: Chart title (default: None)
            has_legend: Whether to show legend (default: True)
            legend_position: Legend position (default: "right")
            **kwargs: "value_column" and "chart_title" are accepted as aliases;
                any other unknown parameter triggers a warning

        Returns:
            The created chart object

        Raises:
            ValueError: If neither data nor categories/values are provided
        """
        from easypptx.chart import Chart, extract_chart_series

        if value_columns is None and "value_column" in kwargs:
            value_columns = kwargs.pop("value_column")
        if title is None and "chart_title" in kwargs:
            title = kwargs.pop("chart_title")
        warn_ignored_kwargs("Slide.add_chart", kwargs)

        # A ChartStyle fills any chart options the caller left unset
        if style is not None:
            sk = style.to_kwargs()
            chart_type = chart_type if chart_type is not None else sk.get("chart_type")
            has_legend = has_legend if has_legend is not None else sk.get("has_legend")
            legend_position = legend_position if legend_position is not None else sk.get("legend_position")
        chart_type = chart_type if chart_type is not None else "column"
        has_legend = has_legend if has_legend is not None else True
        legend_position = legend_position if legend_position is not None else "right"

        series: dict[str, list] | None = None
        if categories is not None and values is not None:
            series = {"Series 1": values}
        else:
            if data is None:
                raise ValueError("Provide either 'data' or both 'categories' and 'values'")
            categories, series = extract_chart_series(data, category_column, value_columns)

        return Chart(self).add(
            chart_type=chart_type,
            categories=categories,
            series=series,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            has_legend=has_legend,
            legend_position=legend_position,
        )

    @property
    def notes(self) -> str:
        """The slide's speaker notes (empty string if none exist)."""
        if not self.pptx_slide.has_notes_slide:
            return ""
        frame = self.pptx_slide.notes_slide.notes_text_frame
        return frame.text if frame is not None else ""

    @notes.setter
    def notes(self, text: str) -> None:
        """Set the slide's speaker notes."""
        frame = self.pptx_slide.notes_slide.notes_text_frame
        if frame is None:  # pragma: no cover - malformed notes master
            raise ValueError("This slide's notes layout has no notes text frame")
        frame.text = text

    def add_bullets(
        self,
        items: list[str | tuple[str, int]],
        x: PositionType = 5,
        y: PositionType = 20,
        width: PositionType = 90,
        height: PositionType = 70,
        font_size: int | None = None,
        font_name: str | None = None,
        font_bold: bool | None = None,
        font_italic: bool | None = None,
        color: str | tuple[int, int, int] | None = None,
        align: str | None = None,
        bullet: bool = True,
        style: TextStyle | None = None,
        **kwargs: Any,
    ) -> PPTXShape:
        """Add a bulleted (or plain stacked) list of paragraphs to the slide.

        Args:
            items: List entries; each item is a string (level 0) or a
                (text, level) tuple for nested bullets (level 0-4)
            x: X position as percent or in_() length (default: 5)
            y: Y position as percent or in_() length (default: 20)
            width: Width as percent or in_() length (default: 90)
            height: Height as percent or in_() length (default: 70)
            font_size: Font size in points (default: 18)
            font_name: Font name (default: "Meiryo")
            color: Text color name or RGB tuple (default: None)
            align: Text alignment "left"/"center"/"right" (default: "left")
            bullet: Draw bullet characters; False gives plain stacked paragraphs (default: True)
            **kwargs: Unknown parameters trigger a warning

        Returns:
            The created text box shape

        Examples:
            ```python
            slide.add_bullets([
                "Revenue up 12%",
                ("EMEA strongest region", 1),
                "Costs down 3%",
            ])
            ```
        """
        warn_ignored_kwargs("Slide.add_bullets", kwargs)

        if style is not None:
            sk = style.to_kwargs()
            font_size = font_size if font_size is not None else sk.get("font_size")
            font_name = font_name if font_name is not None else sk.get("font_name")
            font_bold = font_bold if font_bold is not None else sk.get("font_bold")
            font_italic = font_italic if font_italic is not None else sk.get("font_italic")
            align = align if align is not None else sk.get("align")
            color = color if color is not None else sk.get("color")

        font_size_val = self._template_value("text", "font_size", font_size, 18)
        font_name_val = self._template_value("text", "font_name", font_name, DEFAULT_FONT)
        align_val = self._template_value("text", "align", align, "left")
        color_rgb = resolve_color(normalize_color(self._template_value("text", "color", color, None)))

        x_inches = self._convert_position(x, self._slide_width)
        y_inches = self._convert_position(y, self._slide_height)
        width_inches = self._convert_position(width, self._slide_width)
        height_inches = self._convert_position(height, self._slide_height)

        text_box = self.pptx_slide.shapes.add_textbox(
            Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(items):
            text, level = item if isinstance(item, tuple) else (item, 0)
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = str(text)
            paragraph.level = max(0, min(4, int(level)))
            paragraph.font.size = Pt(font_size_val)
            paragraph.font.name = font_name_val
            if font_bold is not None:
                paragraph.font.bold = font_bold
            if font_italic is not None:
                paragraph.font.italic = font_italic
            if align_val in ALIGN:
                paragraph.alignment = ALIGN[align_val]
            if color_rgb is not None:
                paragraph.font.color.rgb = color_rgb
            if bullet:
                _set_bullet(paragraph)
            else:
                _remove_bullet(paragraph)

        return text_box

    def add_pyplot(
        self,
        figure: Any,
        x: PositionType = 5,
        y: PositionType = 15,
        width: PositionType = 90,
        height: PositionType = 75,
        dpi: int = 300,
        file_format: str = "png",
        border: bool = False,
        border_color: str = "black",
        border_width: int = 1,
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
        **kwargs: Any,
    ) -> PPTXShape:
        """Add a matplotlib (or seaborn) figure to the slide.

        Args:
            figure: Matplotlib figure object (plt.figure(), plt.gcf(), or a
                seaborn plot's .figure)
            x: X position as percent or in_() length (default: 5)
            y: Y position as percent or in_() length (default: 15)
            width: Width as percent or in_() length (default: 90)
            height: Height as percent or in_() length (default: 75)
            dpi: Resolution for the figure (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            border: Whether to draw a border around the image (default: False)
            border_color: Border color name or RGB tuple (default: "black")
            border_width: Border width in points (default: 1)
            shadow: Whether to apply a drop shadow (default: False)
            maintain_aspect_ratio: Whether to keep the figure's aspect ratio (default: True)
            **kwargs: Unknown parameters trigger a warning

        Returns:
            The created picture shape
        """
        from easypptx.pyplot import Pyplot

        warn_ignored_kwargs("Slide.add_pyplot", kwargs)

        return Pyplot.add(
            slide=self,
            figure=figure,
            position={"x": x, "y": y, "width": width, "height": height},
            dpi=dpi,
            file_format=file_format,
            style={
                "border": border,
                "border_color": border_color,
                "border_width": border_width,
                "shadow": shadow,
                "maintain_aspect_ratio": maintain_aspect_ratio,
            },
        )

    def add_multiple_objects(
        self,
        objects_data: list[dict[str, Any]],
        layout: str = "grid",
        padding_percent: float = 5.0,
        start_x: PositionType = "5%",
        start_y: PositionType = "5%",
        width: PositionType = "90%",
        height: PositionType = "90%",
        **kwargs: Any,
    ) -> list[PPTXShape]:
        """Add multiple objects to the slide with automatic alignment.

        Args:
            objects_data: List of dictionaries containing object data
                Each dict should have 'type' ('text', 'image', or 'shape') and type-specific parameters
            layout: Layout type ('grid', 'horizontal', 'vertical')
            padding_percent: Padding between objects as percentage of container
            start_x: Starting X position of container in inches or percentage
            start_y: Starting Y position of container in inches or percentage
            width: Width of container in inches or percentage
            height: Height of container in inches or percentage
            **kwargs: Unknown parameters trigger a warning

        Returns:
            List of created shape objects
        """
        warn_ignored_kwargs("Slide.add_multiple_objects", kwargs)

        # Convert container position and size to inches
        container_x = self._convert_position(start_x, self._slide_width)
        container_y = self._convert_position(start_y, self._slide_height)
        container_width = self._convert_position(width, self._slide_width)
        container_height = self._convert_position(height, self._slide_height)

        padding = padding_percent / 100.0
        num_objects = len(objects_data)

        if layout == "horizontal":
            cols, rows = num_objects, 1
        elif layout == "vertical":
            cols, rows = 1, num_objects
        else:  # Default to grid
            cols = math.ceil(math.sqrt(num_objects))
            rows = math.ceil(num_objects / cols)

        cell_width = container_width / cols
        cell_height = container_height / rows
        obj_width = Inches(cell_width * (1 - padding))
        obj_height = Inches(cell_height * (1 - padding))

        created_objects = []

        for i, obj_data in enumerate(objects_data):
            col = i % cols
            row = i // cols
            obj_x = Inches(container_x + (col * cell_width) + (cell_width * padding / 2))
            obj_y = Inches(container_y + (row * cell_height) + (cell_height * padding / 2))

            obj_type = obj_data.get("type", "text")

            if obj_type == "text":
                obj = self.add_text(
                    text=obj_data.get("text", ""),
                    x=obj_x,
                    y=obj_y,
                    width=obj_width,
                    height=obj_height,
                    font_size=obj_data.get("font_size", 18),
                    font_bold=obj_data.get("font_bold", False),
                    font_italic=obj_data.get("font_italic", False),
                    font_name=obj_data.get("font_name", DEFAULT_FONT),
                    align=obj_data.get("align", "center"),
                    vertical=obj_data.get("vertical", "middle"),
                    color=obj_data.get("color", "black"),
                )
            elif obj_type == "image":
                obj = self.add_image(
                    image_path=obj_data.get("image_path", ""), x=obj_x, y=obj_y, width=obj_width, height=obj_height
                )
            elif obj_type == "shape":
                obj = self.add_shape(
                    shape_type=obj_data.get("shape_type", MSO_SHAPE.RECTANGLE),
                    x=obj_x,
                    y=obj_y,
                    width=obj_width,
                    height=obj_height,
                    fill_color=obj_data.get("fill_color", None),
                )
            else:
                raise ValueError(f"Unknown object type: {obj_type!r} (use 'text', 'image', or 'shape')")

            created_objects.append(obj)

        return created_objects

    @property
    def title(self) -> str | None:
        """Get the slide title.

        Returns:
            The slide title if it exists, None otherwise
        """
        if self.pptx_slide.shapes.title:
            return self.pptx_slide.shapes.title.text
        return None

    @title.setter
    def title(self, value: str) -> None:
        """Set the slide title.

        Args:
            value: The title text
        """
        if self.pptx_slide.shapes.title:
            self.pptx_slide.shapes.title.text = value

    def set_background_color(self, color: str | tuple[int, int, int], **kwargs: Any) -> None:
        """Set the background color of the slide.

        Args:
            color: Background color as string name from COLORS dict or RGB tuple
            **kwargs: Unknown parameters trigger a warning

        Raises:
            ValueError: If the color name is unknown or the value is not a color
        """
        warn_ignored_kwargs("Slide.set_background_color", kwargs)

        background = self.pptx_slide.background
        background._element.attrib["bwMode"] = "auto"

        fill = background.fill
        fill.solid()

        if isinstance(color, str):
            if color not in COLORS:
                raise ValueError(f"Color '{color}' not found in predefined colors")
            fill.fore_color.rgb = COLORS[color]
        elif isinstance(color, tuple) and len(color) == 3:
            fill.fore_color.rgb = RGBColor(*color)
        else:
            raise ValueError("Color must be a string name or RGB tuple")


def _bullet_props(paragraph: Any) -> Any:
    """Get the paragraph-properties XML element, clearing existing bullet settings."""
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    return pPr


def _insert_bullet_element(pPr: Any, element: Any) -> None:
    """Insert a bullet element in schema order (before a:tabLst/a:defRPr/a:extLst)."""
    from pptx.oxml.ns import qn

    for anchor_tag in ("a:tabLst", "a:defRPr", "a:extLst"):
        anchor = pPr.find(qn(anchor_tag))
        if anchor is not None:
            anchor.addprevious(element)
            return
    pPr.append(element)


def _set_bullet(paragraph: Any, char: str = "\u2022") -> None:
    """Mark a paragraph as a bulleted list entry."""
    from pptx.oxml.ns import qn

    pPr = _bullet_props(paragraph)
    _insert_bullet_element(pPr, pPr.makeelement(qn("a:buChar"), {"char": char}))


def _remove_bullet(paragraph: Any) -> None:
    """Mark a paragraph as having no bullet."""
    from pptx.oxml.ns import qn

    pPr = _bullet_props(paragraph)
    _insert_bullet_element(pPr, pPr.makeelement(qn("a:buNone"), {}))


# Backward-compatible export: PositionType historically lived in this module
__all__ = ["PositionType", "Slide"]
