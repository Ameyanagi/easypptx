"""Tests for 0.8.0: data adapter, chart routing, native styling, accessor, text fit."""

import warnings

import numpy as np
import pandas as pd
import polars as pl
import pytest

from easypptx import Presentation, Theme
from easypptx.data import normalize_chart_data, normalize_table_rows


class TestDataAdapter:
    def test_pandas_series(self):
        s = pd.Series([1, 2], index=["a", "b"], name="Rev")
        cats, series = normalize_chart_data(s)
        assert cats == ["a", "b"]
        assert series == {"Rev": [1, 2]}

    def test_dict_of_sequences(self):
        cats, series = normalize_chart_data({"Rev": [1, 2], "Cost": [3, 4]}, categories=["Q1", "Q2"])
        assert cats == ["Q1", "Q2"]
        assert list(series) == ["Rev", "Cost"]

    def test_numpy_2d_with_labels(self):
        cats, series = normalize_chart_data(np.array([[1, 3], [2, 4]]), columns=["Rev", "Cost"], categories=["a", "b"])
        assert series == {"Rev": [1, 2], "Cost": [3, 4]}

    def test_numpy_column_count_mismatch(self):
        with pytest.raises(ValueError, match="columns has"):
            normalize_chart_data(np.array([[1, 2]]), columns=["only-one"])

    def test_polars_dataframe(self):
        df = pl.DataFrame({"Q": ["a", "b"], "R": [1, 2], "C": [3, 4]})
        cats, series = normalize_chart_data(df, value_columns=["R", "C"])
        assert cats == ["a", "b"]
        assert series == {"R": [1, 2], "C": [3, 4]}

    def test_unsupported_type_raises(self):
        with pytest.raises(ValueError, match="Unsupported chart data type"):
            normalize_chart_data(42)

    def test_table_rows_from_all_shapes(self):
        assert normalize_table_rows(pl.DataFrame({"A": [1], "B": [2]})) == [["A", "B"], [1, 2]]
        assert normalize_table_rows({"A": [1], "B": [2]}) == [["A", "B"], [1, 2]]
        assert normalize_table_rows(np.array([[7]]))[0] == ["Column 1"]
        assert normalize_table_rows(pd.Series([5], index=["r"], name="V")) == [["", "V"], ["r", 5]]


class TestChartRouting:
    def test_native_stays_native(self):
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1]], chart_type="column")
        assert hasattr(chart, "plots")  # a native chart object

    def test_pyplot_only_type_routes_to_image(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_chart(data=np.random.rand(3, 4), chart_type="heatmap")
        assert not hasattr(shape, "plots")  # a picture shape

    def test_forced_pyplot_backend(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_chart(data=[["Q", "V"], ["a", 1], ["b", 2]], chart_type="bar", backend="pyplot")
        assert not hasattr(shape, "plots")

    def test_native_backend_rejects_pyplot_type(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="no native PowerPoint equivalent"):
            slide.add_chart(data=[["Q", "V"], ["a", 1]], chart_type="heatmap", backend="native")

    def test_unknown_backend_rejected(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="Unknown backend"):
            slide.add_chart(data=[["Q", "V"], ["a", 1]], chart_type="column", backend="plotly")

    def test_box_violin_histogram_render(self):
        pres = Presentation()
        slide = pres.add_slide()
        data = {"A": [1.0, 2.0, 3.0, 4.0], "B": [2.0, 3.0, 4.0, 5.0]}
        for kind in ("box", "violin", "histogram"):
            shape = slide.add_chart(data=data, chart_type=kind, height=25)
            assert shape is not None


class TestNativeStyling:
    def test_data_labels_and_number_format(self):
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1000], ["b", 2000]], show_values=True, number_format="#,##0")
        assert chart.plots[0].has_data_labels
        assert chart.plots[0].data_labels.number_format == "#,##0"

    def test_axis_titles_and_limits(self):
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(
            data=[["Q", "V"], ["a", 1], ["b", 2]], x_title="Quarter", y_title="USD", y_min=0, y_max=10
        )
        assert chart.value_axis.axis_title.text_frame.text == "USD"
        assert chart.category_axis.axis_title.text_frame.text == "Quarter"
        assert chart.value_axis.minimum_scale == 0
        assert chart.value_axis.maximum_scale == 10

    def test_theme_palette_colors_series(self):
        pres = Presentation(theme="dark")
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "R", "C"], ["a", 1, 2]], value_columns=["R", "C"])
        s0 = chart.plots[0].series[0]
        assert s0.format.fill.fore_color.rgb is not None

    def test_pie_with_axis_options_warns_not_crashes(self):
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1]], chart_type="pie", y_title="ignored")
        assert chart is not None


class TestTableFormatting:
    def test_number_format_per_column(self):
        pres = Presentation()
        slide = pres.add_slide()
        table = slide.add_table([["Item", "Sales"], ["A", 1234.5]], number_format={"Sales": "{:,.1f}"})
        assert table.cell(1, 1).text == "1,234.5"
        assert table.cell(1, 0).text == "A"

    def test_shading_uses_raw_values(self):
        pres = Presentation()
        slide = pres.add_slide()
        table = slide.add_table(
            [["Item", "V"], ["a", 0], ["b", 100]],
            number_format={"V": "{:,.0f}"},
            shade_columns=["V"],
            shade_color="blue",
        )
        low = table.cell(1, 1).fill.fore_color.rgb
        high = table.cell(2, 1).fill.fore_color.rgb
        assert low == (255, 255, 255)  # min -> white
        assert high != low

    def test_unknown_shade_column_raises(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="shade column"):
            slide.add_table([["A"], [1]], shade_columns=["nope"])


class TestPandasAccessor:
    def test_accessor_table_and_chart(self):
        import easypptx

        assert easypptx.register_pandas_accessor()
        pres = Presentation()
        slide = pres.add_slide()
        df = pd.DataFrame({"Region": ["E", "W"], "Sales": [1, 2]})
        table = df.pptx.table(slide, y=20, height=25)
        chart = df.pptx.chart(slide, kind="column", y=50, height=40)
        assert table.cell(0, 0).text == "Region"
        assert hasattr(chart, "plots")


class TestTextFit:
    def test_long_text_shrinks(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_text("word " * 300, x=5, y=5, width=40, height=15, font_size=24)
        assert shape.text_frame.paragraphs[0].font.size.pt < 24

    def test_short_text_keeps_size(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_text("Short", x=5, y=5, width=40, height=15, font_size=24)
        assert shape.text_frame.paragraphs[0].font.size.pt == 24

    def test_fit_none_leaves_size(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_text("word " * 300, x=5, y=5, width=40, height=10, font_size=24, fit="none")
        assert shape.text_frame.paragraphs[0].font.size.pt == 24

    def test_fit_resize_mode(self):
        from pptx.enum.text import MSO_AUTO_SIZE

        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_text("grows", fit="resize")
        assert shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def test_invalid_fit_mode_raises(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="Unknown fit mode"):
            slide.add_text("x", fit="stretch")

    def test_bullets_shrink_to_fit(self):
        pres = Presentation()
        slide = pres.add_slide()
        items = [f"bullet line number {i} with plenty of words attached" for i in range(30)]
        shape = slide.add_bullets(items, x=5, y=5, width=40, height=30, font_size=20)
        assert shape.text_frame.paragraphs[0].font.size.pt < 20

    def test_cjk_estimated_wider(self):
        from easypptx.textfit import estimate_lines

        latin = "hello " * 20
        cjk = "こんにちは" * 24  # same char count as latin
        assert estimate_lines(cjk, 18, 5.0) > estimate_lines(latin, 18, 5.0)


class TestThemePalette:
    def test_custom_theme_palette_flows_to_charts(self):
        theme = Theme(palette=[(10, 20, 30)])
        pres = Presentation(theme=theme)
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1]])
        assert chart.plots[0].series[0].format.fill.fore_color.rgb == (10, 20, 30)

    def test_explicit_palette_beats_theme(self):
        pres = Presentation(theme="dark")
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1]], palette=[(1, 2, 3)])
        assert chart.plots[0].series[0].format.fill.fore_color.rgb == (1, 2, 3)


class TestNoWarningsEndToEnd:
    def test_mixed_deck_warning_free(self, tmp_path):
        with warnings.catch_warnings():
            warnings.simplefilter("error")
            pres = Presentation(theme="corporate")
            slide, grid = pres.add_grid_slide(rows=2, cols=2, title="Data")
            df = pd.DataFrame({"Q": ["a", "b"], "R": [1.0, 2.0]})
            grid[0, 0].add_table(data=df)
            slide.add_chart(data=df, y=55, height=40, width=45, show_values=True)
            slide.add_text("note", x=55, y=55, width=40, height=10)
            pres.save(tmp_path / "deck.pptx")


class TestDocsAgentRegressions:
    def test_categories_forwarded_with_data(self):
        """Regression: categories= used to be dropped when combined with data=."""
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(data={"Rev": [1, 2, 3]}, categories=["Q1", "Q2", "Q3"], chart_type="column")
        assert list(chart.plots[0].categories) == ["Q1", "Q2", "Q3"]

    def test_pie_axis_options_warn(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.warns(UserWarning, match="pie"):
            slide.add_chart(data=[["Q", "V"], ["a", 1]], chart_type="pie", y_title="nope")


class TestCodexReviewRegressions:
    """Regressions from the independent 0.8.0 review."""

    def test_headerless_table_formats_and_shades_all_rows(self):
        pres = Presentation()
        slide = pres.add_slide()
        table = slide.add_table(
            [[1, 10.0], [2, 20.0]],
            has_header=False,
            number_format={1: "{:,.1f}"},
            shade_columns=[1],
        )
        assert table.cell(0, 1).text == "10.0"
        assert table.cell(0, 1).fill.fore_color.rgb == (255, 255, 255)  # min shaded white

    def test_integer_column_labels_prefer_names(self):
        from easypptx.table import apply_number_format

        rows = [[1, 2], ["a", 5.0], ["b", 7.0]]  # header row has integer names 1 and 2
        out = apply_number_format(rows, {1: "{:,.1f}"})
        # Key 1 matches the column NAMED 1 (index 0, non-numeric so unformatted)
        # and must not also positionally format index 1
        assert out[1] == ["a", 5.0]

    def test_nan_and_numpy_scalars_in_shading(self):
        pres = Presentation()
        slide = pres.add_slide()
        table = slide.add_table(
            [["V"], [np.float64(1.0)], [float("nan")], [np.float64(3.0)]],
            shade_columns=["V"],
        )
        # numpy scalars shade; NaN is skipped without crashing
        assert table.cell(1, 0).fill.fore_color.rgb == (255, 255, 255)
        assert table.cell(3, 0).fill.fore_color.rgb != (255, 255, 255)

    def test_fit_never_enlarges_small_fonts(self):
        from easypptx.textfit import fit_font_size

        assert fit_font_size(["tiny"], 5.0, 5.0, font_size=6) == 6

    def test_multiline_text_all_paragraphs_formatted(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_text("line one\nline two", font_size=30, fit="none")
        sizes = [p.font.size.pt for p in shape.text_frame.paragraphs]
        assert sizes == [30, 30]

    def test_pyplot_multi_series_pie_rejected(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="single series"):
            slide.add_chart(data={"A": [1, 2], "B": [3, 4]}, chart_type="pie", backend="pyplot", categories=["x", "y"])

    def test_structured_numpy_array_rejected(self):
        arr = np.array([(1, 2.0)], dtype=[("a", "i4"), ("b", "f4")])
        with pytest.raises(ValueError, match="Structured numpy arrays"):
            normalize_chart_data(arr)

    def test_polars_out_of_range_int_column(self):
        with pytest.raises(ValueError, match="not found"):
            normalize_chart_data(pl.DataFrame({"A": [1]}), value_columns=[5])

    def test_empty_pandas_frame(self):
        assert normalize_chart_data(pd.DataFrame()) == ([], {})

    def test_foreign_pptx_accessor_not_overridden(self):
        import easypptx

        assert easypptx.register_pandas_accessor()  # ours registers/detects fine
        # Simulate a foreign accessor
        original = pd.DataFrame.pptx
        try:

            class Foreign:
                pass

            pd.DataFrame.pptx = Foreign()
            with pytest.warns(UserWarning, match="another library"):
                assert not easypptx.register_pandas_accessor()
        finally:
            pd.DataFrame.pptx = original


class TestGreptileReviewRegressions:
    """Regressions from the Greptile PR-review comments (PRs 3-5)."""

    def test_append_skips_spanned_cells_and_uses_absolute_coords(self):
        from pptx.util import Emu

        pres = Presentation()
        slide = pres.add_slide()
        grid = pres.add_grid(slide=slide, x="50%", y="50%", width="50%", height="50%", rows=2, cols=2, padding=0.0)
        grid[0, :]  # span the top row -> (0,1) is covered
        grid.append(lambda **kw: slide.add_text("a", **kw))
        shape = grid.append(lambda **kw: slide.add_text("b", **kw))
        assert grid.cells[0][1].content is None  # spanned cell untouched
        assert grid.cells[1][0].content is not None
        # Placed with absolute slide coordinates (inside the half-slide grid)
        assert shape.left >= Emu(int(12192768 * 0.49))

    def test_expand_preserves_merged_geometry(self):
        from easypptx.positioning import parse_percent

        pres = Presentation()
        slide = pres.add_slide()
        grid = pres.add_grid(slide=slide, rows=2, cols=2, padding=0.0)
        grid[0, :]
        before = parse_percent(grid.cells[0][0].width)
        grid._expand_grid(add_rows=1)
        assert parse_percent(grid.cells[0][0].width) == pytest.approx(before, abs=0.1)

    def test_pie_axis_warning_does_not_skip_palette(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.warns(UserWarning, match="pie"):
            chart = slide.add_chart(
                data=[["Q", "V"], ["a", 1]], chart_type="pie", y_title="ignored", palette=[(9, 9, 9)]
            )
        point_or_series = chart.plots[0].series[0]
        assert point_or_series.format.fill.fore_color.rgb == (9, 9, 9)

    def test_list_data_out_of_range_int_column(self):
        with pytest.raises(ValueError, match="out of range"):
            normalize_chart_data([["A", "B"], [1, 2]], value_columns=[9])
        with pytest.raises(ValueError, match="out of range"):
            normalize_chart_data([["A", "B"], [1, 2]], category_column=7)


class TestThemedFigures:
    """Visual-QA regressions: figures must be legible on themed decks."""

    def test_native_chart_text_uses_theme_color(self):
        pres = Presentation(theme="dark")
        slide = pres.add_slide()
        chart = slide.add_chart(
            data=[["Q", "R", "C"], ["a", 1, 2], ["b", 2, 3]],
            value_columns=["R", "C"],
            show_values=True,
            y_title="USD",
        )
        white = (255, 255, 255)
        assert chart.value_axis.tick_labels.font.color.rgb == white
        assert chart.legend.font.color.rgb == white
        assert chart.plots[0].data_labels.font.color.rgb == white

    def test_pyplot_chart_transparent_on_theme(self):
        from PIL import Image as PILImage

        pres = Presentation(theme="dark")
        slide = pres.add_slide()
        shape = slide.add_chart(data={"A": [1.0, 2.0, 3.0]}, chart_type="histogram", categories=["x", "y", "z"])
        import io

        img = PILImage.open(io.BytesIO(shape.image.blob))
        assert img.mode == "RGBA"
        corner_alpha = img.getpixel((0, 0))[3]
        assert corner_alpha == 0  # transparent background

    def test_explicit_font_color_beats_theme(self):
        pres = Presentation(theme="dark")
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "V"], ["a", 1]], font_color=(1, 2, 3))
        assert chart.value_axis.tick_labels.font.color.rgb == (1, 2, 3)

    def test_grid_cell_table_supports_formatting(self):
        pres = Presentation()
        slide, grid = pres.add_grid_slide(rows=1, cols=1, title="T")
        table = grid[0, 0].add_table(
            data=[["Item", "Sales"], ["A", 1234.5], ["B", 42.0]],
            number_format={"Sales": "{:,.1f}"},
            shade_columns=["Sales"],
        )
        assert table.cell(1, 1).text == "1,234.5"
        assert table.cell(2, 1).fill.fore_color.rgb == (255, 255, 255)


class TestLegendLayout:
    def test_legend_does_not_overlap_plot(self):
        """PowerPoint visual QA: the legend must reserve its own layout space."""
        pres = Presentation()
        slide = pres.add_slide()
        chart = slide.add_chart(data=[["Q", "R", "C"], ["a", 1, 2]], value_columns=["R", "C"])
        assert chart.legend.include_in_layout is False
