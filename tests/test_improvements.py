"""Behavior tests for the 0.6.0 correctness and API improvements."""

import warnings

import pandas as pd
import pytest
from pptx.util import Emu

from easypptx import Presentation
from easypptx.positioning import pct, shift_band, to_percent

SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches (16:9)
SLIDE_WIDTH_EMU = 12192768  # 13.33 inches (16:9)


def emu_to_height_percent(emu: int) -> float:
    """Convert an EMU y-coordinate to a percentage of the 16:9 slide height."""
    return emu / SLIDE_HEIGHT_EMU * 100


class TestPositioningHelpers:
    def test_shift_band_percent(self):
        y, height = shift_band("10%", "80%", "10%", SLIDE_HEIGHT_EMU)
        assert y == "20.00%"
        assert height == "70.00%"

    def test_shift_band_inches(self):
        from easypptx import in_

        # 0.75 inches is 10% of a 7.5-inch-high slide
        y, height = shift_band(in_(0.75), in_(6.0), in_(0.75), SLIDE_HEIGHT_EMU)
        assert y == "20.00%"
        assert height == "70.00%"

    def test_to_percent_and_pct(self):
        from easypptx import in_

        assert to_percent("25%", SLIDE_HEIGHT_EMU) == 25.0
        assert to_percent(7.5, SLIDE_HEIGHT_EMU) == 7.5  # bare numbers are percent
        assert to_percent(in_(7.5), SLIDE_HEIGHT_EMU) == pytest.approx(100.0)
        assert pct(12.5) == "12.50%"


class TestContentPaddingAppliedOnce:
    def test_grid_slide_content_y_padding_applied_once(self):
        """Regression: content_y_padding used to shift and shrink the grid twice."""
        pres = Presentation()
        _slide, grid = pres.add_grid_slide(rows=2, cols=2, title="T", content_y_padding="5%")

        # title_height defaults to 10%; padding of 5% applied exactly once
        assert to_percent(grid.y, SLIDE_HEIGHT_EMU) == pytest.approx(15.0, abs=0.05)
        assert to_percent(grid.height, SLIDE_HEIGHT_EMU) == pytest.approx(85.0, abs=0.05)


class TestSlideContentMethods:
    def test_slide_add_table(self):
        pres = Presentation()
        slide = pres.add_slide()
        table = slide.add_table([["Name", "Value"], ["A", 1]], x="10%", y="20%", has_header=True)
        assert table.cell(0, 0).text == "Name"

    def test_slide_add_chart_with_value_column_alias(self):
        pres = Presentation()
        slide = pres.add_slide()
        data = [["Region", "Sales"], ["East", 10], ["West", 20]]
        chart = slide.add_chart(data=data, chart_type="bar", category_column="Region", value_column="Sales")
        assert chart is not None

    def test_slide_add_shape_accepts_string_type(self):
        pres = Presentation()
        slide = pres.add_slide()
        shape = slide.add_shape(shape_type="ROUNDED_RECTANGLE", x="10%", y="10%", width="20%", height="10%")
        assert shape is not None

    def test_slide_add_shape_rejects_unknown_string(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.raises(ValueError, match="Unknown shape type"):
            slide.add_shape(shape_type="NOT_A_SHAPE")


class TestLoudFailures:
    def test_unknown_kwargs_warn(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.warns(UserWarning, match="font_szie"):
            slide.add_text("typo", font_szie=24)

    def test_vertical_align_alias_accepted_silently(self):
        pres = Presentation()
        slide = pres.add_slide()
        with warnings.catch_warnings():
            warnings.simplefilter("error")
            slide.add_text("aliased", vertical_align="middle")

    def test_out_of_range_percent_warns(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.warns(UserWarning, match="clamped"):
            slide.add_text("off-slide", x="150%")

    def test_bad_template_toml_raises(self):
        pres = Presentation()
        with pytest.raises(FileNotFoundError):
            pres.add_slide(template_toml="no_such_template.toml")

    def test_reference_swap_after_slides_raises(self):
        import easypptx as ep

        pres = Presentation()
        pres.add_slide()
        other_ref = str(__import__("pathlib").Path(ep.__file__).parent / "reference_4x3.pptx")
        pres.template_manager.register("swap_template", {"bg_color": "white"})
        pres.template_manager.template_references["swap_template"] = other_ref
        with pytest.raises(ValueError, match="already contains slides"):
            pres.add_slide_from_template("swap_template")


class TestSlideWrapperCaching:
    def test_slides_property_preserves_user_data(self):
        pres = Presentation()
        slide = pres.add_slide()
        slide.user_data["key"] = "value"
        assert pres.slides[0] is slide
        assert pres.slides[0].user_data["key"] == "value"


class TestTableFixes:
    def test_from_dataframe_include_index(self):
        """Regression: include_index=True used to raise a column-count error."""
        from easypptx.table import Table

        pres = Presentation()
        slide = pres.add_slide()
        df = pd.DataFrame({"A": [1, 2], "B": [3, 4]})
        table = Table(slide).from_dataframe(df, include_index=True)
        assert table.cell(0, 0).text == ""
        assert table.cell(1, 0).text == "0"

    def test_has_header_alias(self):
        pres = Presentation()
        slide = pres.add_slide()
        from easypptx.table import Table

        table = Table(slide).add([["H", "I"], [1, 2]], has_header=False)
        assert table.cell(0, 0).text_frame.paragraphs[0].font.bold is None


class TestGridGeometry:
    def test_grid_cell_content_positioned_absolutely(self):
        """Content in a half-slide grid must land inside the grid, not at cell-relative coords."""
        pres = Presentation()
        slide = pres.add_slide()
        grid = pres.add_grid(slide=slide, x="50%", y="50%", width="50%", height="50%", rows=1, cols=1, padding=0.0)
        shape = grid[0, 0].add_text("positioned")
        # The cell starts at the grid origin: 50% of slide width/height
        assert shape.left >= Emu(int(SLIDE_WIDTH_EMU * 0.49))
        assert shape.top >= Emu(int(SLIDE_HEIGHT_EMU * 0.49))

    def test_autogrid_positions_position_aware_content(self):
        """Content functions accepting x/y/width/height receive the cell position."""
        received = []

        pres = Presentation()
        slide = pres.add_slide()

        def make_func(idx):
            def func(**kwargs):
                received.append(kwargs)
                return slide.add_text(f"item {idx}", **kwargs)

            return func

        pres.add_autogrid(slide, [make_func(0), make_func(1)], rows=1, cols=2)
        assert len(received) == 2
        assert received[0]["x"] != received[1]["x"]


class TestDeprecations:
    def test_pass_through_methods_warn(self):
        pres = Presentation()
        slide = pres.add_slide()
        with pytest.warns(DeprecationWarning, match="Slide.add_text"):
            pres.add_text(slide, "deprecated path")


class TestCodexReviewRegressions:
    """Regressions surfaced by the independent post-refactor review."""

    def test_chart_extraction_with_integer_column_labels(self):
        from easypptx.chart import extract_categories_values

        df = pd.DataFrame([["A", 1], ["B", 2]])  # RangeIndex columns: 0, 1
        categories, values = extract_categories_values(df)
        assert categories == ["A", "B"]
        assert values == [1, 2]

    def test_same_reference_pptx_does_not_raise_after_slides(self):
        import pathlib

        import easypptx as ep

        # The default 16:9 presentation already uses this reference file
        same_ref = str(pathlib.Path(ep.__file__).parent / "reference_16x9.pptx")
        pres = Presentation()
        pres.add_slide()
        pres.template_manager.register("same_ref_template", {"bg_color": "white"})
        pres.template_manager.template_references["same_ref_template"] = same_ref
        slide = pres.add_slide_from_template("same_ref_template")
        assert slide is not None

    def test_grid_pyplot_with_global_defaults_does_not_warn(self):
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        pres = Presentation()
        slide = pres.add_slide()
        grid = pres.add_grid(slide=slide, rows=1, cols=1)
        # Broad global defaults, as produced by generate_default_template
        grid.template_defaults["global"] = {"font_size": 14, "title_align": "center", "chart_type": "bar"}

        fig = plt.figure()
        plt.plot([1, 2], [3, 4])
        with warnings.catch_warnings():
            warnings.simplefilter("error")
            grid.add_pyplot(0, 0, figure=fig)
        plt.close(fig)

    def test_template_cache_respects_requested_name(self, tmp_path):
        import tomli_w

        from easypptx.template import TemplateManager

        path = tmp_path / "mytemplate.toml"
        with open(path, "wb") as f:
            tomli_w.dump({"bg_color": "white"}, f)

        tm = TemplateManager(template_dir=str(tmp_path))
        assert tm.load(str(path), template_name="custom") == "custom"
        # A later plain load must register the filename stem, not return "custom"
        assert tm.load(str(path)) == "mytemplate"
        assert tm.get("mytemplate") is not None
        # And repeated plain loads hit the cache
        assert tm.load(str(path)) == "mytemplate"
