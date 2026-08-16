"""Tests for the fluent Deck builder (0.10.0)."""

import warnings

import pandas as pd
import pytest

from easypptx import Deck, Presentation

DF = pd.DataFrame({"Q": ["Q1", "Q2"], "Rev": [1, 2], "Cost": [3, 4]})


def slide_texts(slide):
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame]


class TestChain:
    def test_full_chain_renders_expected_slides(self):
        pres = (
            Deck(theme="corporate")
            .title_slide("Q3 Review", subtitle="Finance", notes="opening")
            .slide("Highlights", notes="short")
            .bullets(["Revenue up 12%", ("EMEA strongest", 1)])
            .chart(DF, kind="column", value_columns=["Rev", "Cost"])
            .section("Details")
            .slide("Data")
            .table(DF)
            .text("Totals in appendix.")
            .build()
        )
        assert len(pres.slides) == 4
        assert pres.slides[0].notes == "opening"
        assert pres.slides[1].notes == "short"
        assert "Q3 Review" in slide_texts(pres.slides[0])
        assert "Details" in slide_texts(pres.slides[2])

    def test_every_method_returns_the_deck(self):
        deck = Deck()
        assert deck.slide("a") is deck
        assert deck.text("t") is deck
        assert deck.bullets(["b"]) is deck
        assert deck.chart(DF) is deck
        assert deck.table(DF) is deck
        assert deck.notes("n") is deck
        assert deck.tap(lambda s: None) is deck

    def test_save_writes_file(self, tmp_path):
        out = tmp_path / "deck.pptx"
        Deck().slide("One").text("hello").save(out)
        assert out.exists()

    def test_pres_deck_builds_onto_existing_presentation(self):
        pres = Presentation(theme="dark")
        result = pres.deck().slide("A").text("x").build()
        assert result is pres
        assert len(pres.slides) == 1


class TestValidation:
    def test_content_before_slide_raises(self):
        with pytest.raises(RuntimeError, match="Start a slide"):
            Deck().bullets(["a"])

    def test_content_after_title_slide_raises(self):
        with pytest.raises(RuntimeError, match="Start a slide"):
            Deck().title_slide("T").text("body")

    def test_unknown_chart_kind_raises_at_call_site(self):
        with pytest.raises(ValueError, match="starburst"):
            Deck().slide("x").chart(DF, kind="starburst")

    def test_bad_chart_data_raises_at_call_site(self):
        with pytest.raises(ValueError, match="Unsupported chart data type"):
            Deck().slide("x").chart(42)

    def test_missing_image_raises_at_call_site(self):
        with pytest.raises(FileNotFoundError):
            Deck().slide("x").image("nope.png")

    def test_bad_bullets_raise(self):
        with pytest.raises(TypeError):
            Deck().slide("x").bullets("not a list")
        with pytest.raises(TypeError):
            Deck().slide("x").bullets([("text", "level")])

    def test_consumed_deck_cannot_be_reused(self):
        deck = Deck().slide("a")
        deck.build()
        with pytest.raises(RuntimeError, match="already rendered"):
            deck.slide("b")
        with pytest.raises(RuntimeError, match="already rendered"):
            deck.build()


class TestLayout:
    def test_chart_expands_into_leftover_space(self):
        pres = Deck().slide("T").bullets(["one", "two"]).chart(DF).build()
        chart_frame = next(s for s in pres.slides[0].shapes if s.has_chart)
        # The chart should take a large share of the slide height
        assert chart_frame.height > pres.pptx_presentation.slide_height * 0.4

    def test_positioned_block_escapes_auto_layout(self):
        pres = Deck().slide("T").text("pinned", x=60, y=80, width=30, height=10).build()
        shape = next(s for s in pres.slides[0].shapes if "pinned" in s.text_frame.text)
        assert shape.top > pres.pptx_presentation.slide_height * 0.7

    def test_tap_receives_rendered_slide(self):
        seen = []
        Deck().slide("T").text("x").tap(lambda s: seen.append(len(s.shapes))).build()
        assert seen and seen[0] >= 2  # title + text already rendered

    def test_notes_append(self):
        pres = Deck().slide("T", notes="first").notes("second").build()
        assert pres.slides[0].notes == "first\nsecond"

    def test_rendering_is_warning_free(self):
        with warnings.catch_warnings():
            warnings.simplefilter("error")
            (
                Deck(theme="dark")
                .title_slide("T", subtitle="s")
                .slide("Content")
                .bullets(["a", "b", ("c", 1)])
                .chart(DF, value_columns=["Rev", "Cost"], show_values=True)
                .slide("More")
                .table(DF)
                .build()
            )


class TestPagination:
    def test_long_bullets_paginate_with_cont_titles(self):
        deck = Deck(theme="light").slide("Long list")
        deck.bullets([f"Bullet point number {i} with a reasonable amount of words" for i in range(40)])
        pres = deck.build()
        assert len(pres.slides) >= 2
        titles = [slide_texts(s)[0] for s in pres.slides]
        assert titles[0] == "Long list"
        assert all("(cont.)" in t for t in titles[1:])

    def test_notes_only_on_first_page(self):
        deck = Deck().slide("L", notes="only once")
        deck.bullets([f"item {i} with several words to wrap around" for i in range(40)])
        pres = deck.build()
        assert pres.slides[0].notes == "only once"
        assert pres.slides[1].notes == ""

    def test_paginate_false_compresses_instead(self):
        deck = Deck().slide("Dense", paginate=False)
        deck.bullets([f"item {i} with several words attached" for i in range(40)])
        pres = deck.build()
        assert len(pres.slides) == 1

    def test_mixed_blocks_flow_to_next_page(self):
        deck = Deck().slide("Mix")
        deck.bullets([f"filler line {i} with enough words to take space" for i in range(30)])
        deck.table(DF)
        pres = deck.build()
        assert len(pres.slides) >= 2
        # The table must exist on some page
        assert any(any(s.has_table for s in slide.shapes) for slide in pres.slides)


class TestDesignTier:
    """0.10.0 design features: stats, compare, kicker, footer, emphasis."""

    def test_stats_tiles_render(self):
        pres = (
            Deck(theme="corporate")
            .slide("Numbers")
            .stats([("+12%", "Growth", "+3pt"), ("$1.6M", "Revenue"), ("38", "Logos", "-4")])
            .build()
        )
        texts = " ".join(slide_texts(pres.slides[0]))
        assert "+12%" in texts and "GROWTH" in texts and "-4" in texts

    def test_stats_validation(self):
        with pytest.raises(ValueError, match="between 1 and 5"):
            Deck().slide("x").stats([])
        with pytest.raises(TypeError):
            Deck().slide("x").stats(["just a string"])

    def test_compare_cards_render(self):
        pres = Deck(theme="light").slide("Choice").compare(("Now", ["a", "b"]), ("Later", ["c"])).build()
        texts = " ".join(slide_texts(pres.slides[0]))
        assert "Now" in texts and "Later" in texts

    def test_kicker_renders_above_title(self):
        pres = Deck(theme="dark").slide("Results", kicker="Q3 FINANCIALS").text("body").build()
        assert "Q3 FINANCIALS" in slide_texts(pres.slides[0])

    def test_footer_and_page_numbers(self):
        pres = Deck(theme="dark", footer="Acme Corp").slide("A").text("x").slide("B").text("y").build()
        texts = slide_texts(pres.slides[1])
        assert "Acme Corp" in texts
        assert "2 / 2" in texts

    def test_chart_headline_and_emphasize(self):
        pres = (
            Deck(theme="corporate")
            .slide("Msg")
            .chart(DF, value_columns=["Rev", "Cost"], emphasize="Rev", headline="Rev pulls ahead")
            .build()
        )
        assert "Rev pulls ahead" in slide_texts(pres.slides[0])
        chart = next(s for s in pres.slides[0].shapes if s.has_chart).chart
        rev, cost = chart.plots[0].series
        assert rev.format.fill.fore_color.rgb != cost.format.fill.fore_color.rgb

    def test_title_slide_composition(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        pres = Deck(theme="corporate").title_slide("T", subtitle="s").build()
        shapes = pres.slides[0].shapes
        # Panel + edge + accent bar = at least 3 autoshapes
        assert sum(1 for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE) >= 3


class TestCodexBuilderRegressions:
    """Fixes from the independent 0.10.0 builder review."""

    def test_untitled_slides_use_larger_capacity(self):
        items = [f"line {i} with a few words here" for i in range(24)]
        titled = Deck().slide("T")
        titled.bullets(list(items))
        untitled = Deck().slide()
        untitled.bullets(list(items))
        assert len(untitled.build().slides) <= len(titled.build().slides)

    def test_oversized_bullet_moves_to_next_page(self):
        deck = Deck().slide("T")
        deck.text("short intro " * 20)
        deck.bullets(["x " * 600])  # one bullet taller than a page
        pres = deck.build()
        assert len(pres.slides) >= 2

    def test_chart_value_column_alias_validated_eagerly(self):
        with pytest.raises(ValueError, match="not found"):
            Deck().slide("x").chart(DF, value_column="Missing")

    def test_categories_values_length_mismatch_eager(self):
        with pytest.raises(ValueError, match="same length"):
            Deck().slide("x").chart(categories=["a", "b"], values=[1])

    def test_data_mutation_after_call_is_ignored(self):
        rows = [["A", "B"], [1, 2]]
        deck = Deck().slide("T").table(rows)
        rows.extend([[i, i] for i in range(30)])  # mutate after the call
        pres = deck.build()
        table = next(s for s in pres.slides[0].shapes if s.has_table).table
        assert len(table.rows) == 2

    def test_empty_and_ragged_tables_rejected_eagerly(self):
        with pytest.raises(ValueError):
            Deck().slide("x").table([])
        with pytest.raises(ValueError, match="rectangular"):
            Deck().slide("x").table([["A", "B"], [1]])

    def test_tap_runs_after_content(self):
        order = []
        (
            Deck()
            .slide("T")
            .tap(lambda s: order.append(len([sh for sh in s.shapes if sh.has_text_frame])))
            .text("after the tap in chain order")
            .build()
        )
        # The tap saw the completed slide (title + text), not an empty one
        assert order and order[0] >= 2
