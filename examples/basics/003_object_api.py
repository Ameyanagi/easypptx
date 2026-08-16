"""
Example showing the object-related APIs in EasyPPTX.

This example demonstrates how to use the direct object manipulation methods
in the Presentation class to add various elements to slides.
"""

from pathlib import Path

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE

from easypptx import Image, Presentation

# Create a folder for outputs if it doesn't exist
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)

# Create a new presentation
pres = Presentation()

# Create a title slide
title_slide = pres.add_title_slide(title="Object API Examples", subtitle="Direct manipulation of slide elements")

# Slide 1: Text
text_slide = pres.add_content_slide(title="Text Examples")

# Add text directly to the slide
text_slide.add_text(text="This is regular text", x="10%", y="20%", width="80%", height="8%", font_size=18)

text_slide.add_text(
    text="This is bold, italic, blue text",
    x="10%",
    y="30%",
    width="80%",
    height="8%",
    font_size=18,
    font_bold=True,
    font_italic=True,
    color="blue",
)

text_slide.add_text(text="Right-aligned text", x="10%", y="40%", width="80%", height="8%", font_size=18, align="right")

text_slide.add_text(
    text="Middle-aligned text with custom color",
    x="10%",
    y="50%",
    width="80%",
    height="8%",
    font_size=18,
    align="center",
    vertical="middle",
    color=(150, 75, 0),
)

# Slide 2: Shapes
shape_slide = pres.add_content_slide(title="Shape Examples")

# Add various shapes (slide.add_shape does not draw text itself,
# so a label is added on top with slide.add_text)


def add_labeled_shape(slide, shape_type, x, y, width, height, fill_color, label, label_color):
    slide.add_shape(shape_type=shape_type, x=x, y=y, width=width, height=height, fill_color=fill_color)
    slide.add_text(
        text=label,
        x=x,
        y=y,
        width=width,
        height=height,
        align="center",
        vertical="middle",
        color=label_color,
    )


add_labeled_shape(shape_slide, MSO_SHAPE.RECTANGLE, "10%", "20%", "30%", "15%", "blue", "Rectangle", "white")
add_labeled_shape(
    shape_slide, MSO_SHAPE.ROUNDED_RECTANGLE, "60%", "20%", "30%", "15%", "green", "Rounded Rectangle", "white"
)
add_labeled_shape(shape_slide, MSO_SHAPE.OVAL, "10%", "45%", "30%", "15%", "red", "Oval", "white")
add_labeled_shape(shape_slide, MSO_SHAPE.PENTAGON, "60%", "45%", "30%", "15%", "orange", "Pentagon", "white")
add_labeled_shape(shape_slide, MSO_SHAPE.CHEVRON, "10%", "70%", "30%", "15%", "cyan", "Chevron", "black")
add_labeled_shape(shape_slide, MSO_SHAPE.STAR_5_POINT, "60%", "70%", "20%", "15%", "yellow", "Star", "black")

# Slide 3: Table
table_slide = pres.add_content_slide(title="Table Example")

# Create table data
table_data = [
    ["Product", "Q1", "Q2", "Q3", "Q4"],
    ["Product A", 120, 140, 135, 150],
    ["Product B", 85, 90, 95, 110],
    ["Product C", 45, 55, 65, 70],
    ["Product D", 30, 25, 40, 35],
]

# Add table to slide
table_slide.add_table(
    data=table_data,
    x="10%",
    y="20%",
    width="80%",
    height="50%",
    has_header=True,
)

# Slide 4: Chart
chart_slide = pres.add_content_slide(title="Chart Example")

# Create chart data
chart_data = pd.DataFrame({
    "Quarter": ["Q1", "Q2", "Q3", "Q4"],
    "Sales": [120, 145, 160, 180],
    "Expenses": [95, 110, 115, 130],
    "Profit": [25, 35, 45, 50],
})

# Add chart to slide
chart_slide.add_chart(
    data=chart_data,
    chart_type="column",
    x="10%",
    y="20%",
    width="80%",
    height="60%",
    has_legend=True,
    legend_position="bottom",
    category_column="Quarter",
    value_columns=["Sales", "Expenses", "Profit"],
    chart_title="Quarterly Performance",
)

# Slide 5: Image
image_slide = pres.add_content_slide(title="Image Example")

# Create a sample image path
# Check if the output/images directory exists, if not use a placeholder text
image_dir = output_dir / "images"
images = list(image_dir.glob("*.png")) if image_dir.exists() else []
if images:
    sample_image = str(images[0])
else:
    sample_image = None
    image_slide.add_text(
        text="Sample image not found in output/images directory",
        x="10%",
        y="40%",
        width="80%",
        height="10%",
        font_size=16,
        align="center",
    )

# If we have an image, add it to the slide
# (border/shadow styling from the old API was dropped)
if sample_image:
    Image(image_slide).add(
        sample_image,
        x="20%",
        y="20%",
        width="60%",
        maintain_aspect_ratio=True,
    )

# Slide 6: Combination of elements
combined_slide = pres.add_content_slide(title="Combined Elements")

# Add a shape as background panel (line_color styling from the old API was dropped)
combined_slide.add_shape(
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="5%",
    y="15%",
    width="90%",
    height="75%",
    fill_color=(240, 240, 240),
)

# Add title text
combined_slide.add_text(
    text="Dashboard Title",
    x="10%",
    y="20%",
    width="80%",
    height="10%",
    font_size=24,
    font_bold=True,
    align="center",
)

# Add a small table
table_data_small = [
    ["Metric", "Value"],
    ["Total Revenue", "$1,245,000"],
    ["Growth Rate", "+15.2%"],
    ["Conversion", "4.8%"],
]

combined_slide.add_table(data=table_data_small, x="10%", y="35%", width="35%", height="20%", has_header=True)

# Add small chart
chart_data_small = pd.DataFrame({"Month": ["Jan", "Feb", "Mar", "Apr"], "Value": [42, 35, 65, 58]})

combined_slide.add_chart(
    data=chart_data_small,
    chart_type="line",
    x="55%",
    y="35%",
    width="35%",
    height="20%",
    category_column="Month",
    value_columns="Value",
    chart_title="Monthly Trend",
    has_legend=False,
)

# Add labeled KPI boxes at the bottom
for box_x, box_color, box_text in [
    ("10%", "blue", "$1.24M\nTotal Revenue"),
    ("40%", "green", "+15.2%\nGrowth"),
    ("70%", "orange", "4.8%\nConversion"),
]:
    combined_slide.add_shape(
        shape_type=MSO_SHAPE.RECTANGLE, x=box_x, y="65%", width="20%", height="15%", fill_color=box_color
    )
    combined_slide.add_text(
        text=box_text,
        x=box_x,
        y="65%",
        width="20%",
        height="15%",
        font_size=14,
        color="white",
        align="center",
        vertical="middle",
    )

# Add a closing slide
thank_you_slide = pres.add_section_slide(title="Thank You!")

# Save the presentation
output_path = output_dir / "object_api_example.pptx"
pres.save(output_path)
print(f"Presentation saved to {output_path}")
